"""Verify captioned decks — placement + text-presence gate.

v0.2.2: the gate now covers ALL FIVE placement defect classes the auditor
checks, and fails (non-zero exit) on MISSING captions — closing the v0.2.1 gaps
where verify.py exited "clean" while captions were absent and only 2 of 5
placement patterns were checked.

Coordinate checks (all 2D, consistent with apply + audit):
  - Pattern A: no caption overlaps a slide TITLE placeholder.
  - Pattern B: no caption bottom is within FOOTER_CLEARANCE_EMU of the footer.
  - Pattern C: no caption box is geometrically inside its picture
               (SmartArt-icon captions are exempt — they sit inside the
               diagram by design).
  - Pattern D: no caption overlaps a BODY/SUBTITLE placeholder with text.
  - Pattern E: no two captions on the same slide overlap each other.

Audit rows with action='overlay-fullbleed' or 'flagged-no-slot' are treated
as known-missing (intentionally skipped placement), not text-mismatch errors.

Exit codes: 0 = clean; 4 = a placement defect (A/B/C/D/E); 5 = a missing
caption (text-miss); 2 = usage/IO error. The TOTAL line always prints every
count so nothing is hidden behind the code.

Usage:
  python3 verify.py <work_dir>

Reads:
  <work_dir>/captions.json
  <work_dir>/captioned_decks/<deck>_captioned.pptx
  <work_dir>/audit/<deck>_audit.csv  (optional — used to honor known-skip rows)
"""
import sys, os, json, hashlib, glob, csv
from pptx import Presentation
from _oxml_pics import iter_slide_pics, resolve_blob, guess_ext
from _geometry import (
    slide_footer_top, slide_text_obstacle_rects,
    FOOTER_CLEARANCE_EMU, rect_intersect_area, CAPTION_NAME_PREFIXES,
    caption_overflows, band_covers_structural_picture,
)

SA_ICON_PREFIX = 'captioner_sa_icon_'
BAND_PREFIX = 'captioner_capband_'   # caption deliberately in own-picture bottom strip
BG_REPEAT_THRESHOLD = 4              # must match apply_captions --bg-repeat-threshold default


def load_known_skips(audit_dir, deck_stem):
    """Return a set of image_hash values that the apply pass intentionally
    did NOT caption (overlay-fullbleed or flagged-no-slot). Caller treats
    these as known-missing during text-match verification."""
    skips = set()
    p = os.path.join(audit_dir, f"{deck_stem}_audit.csv")
    if not os.path.exists(p):
        return skips
    try:
        with open(p, encoding="utf-8", newline="") as f:
            for row in csv.DictReader(f):
                act = row.get("action", "")
                if ("overlay-fullbleed" in act or "flagged-no-slot" in act
                        or "flagged-self-check" in act
                        or "skipped-decorative-background" in act):
                    h = row.get("image_hash", "")
                    if h:
                        skips.add(h)
    except Exception:
        pass
    return skips


def caption_shapes(slide):
    """Yield (shape_name, (left, top, width, height), text, font_pt) for every
    caption-named shape. font_pt is the first run's size in points (or None) so
    the overflow gate can use the caption's ACTUAL rendered size — icon captions
    are scaled 6–8pt by apply, and an 8pt assumption would false-flag them."""
    for sh in slide.shapes:
        nm = sh.name or ""
        if not any(nm.startswith(p) for p in CAPTION_NAME_PREFIXES):
            continue
        if None in (sh.left, sh.top, sh.width, sh.height):
            continue
        txt = ""
        fpt = None
        try:
            if getattr(sh, "has_text_frame", False):
                txt = sh.text_frame.text
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size is not None:
                            fpt = run.font.size.pt
                            break
                    if fpt is not None:
                        break
        except Exception:
            pass
        yield nm, (int(sh.left), int(sh.top), int(sh.width), int(sh.height)), txt, fpt


def _overlaps_2d(cap_ltwh, obs_ltrb, vfrac=0.15):
    """True if caption (l,t,w,h) overlaps obstacle (L,T,R,B) in BOTH axes,
    vertically by > vfrac of caption height."""
    cl, ct, cw, ch = cap_ltwh
    ol, ot, orr, ob = obs_ltrb
    ix = min(cl + cw, orr) - max(cl, ol)
    iy = min(ct + ch, ob) - max(ct, ot)
    return ix > 0 and iy > vfrac * max(1, ch)


def main():
    if len(sys.argv) != 2:
        print(__doc__); sys.exit(2)
    work = os.path.abspath(sys.argv[1])
    captions_path = os.path.join(work, "captions.json")
    if not os.path.exists(captions_path):
        print(f"ERROR: {captions_path} not found"); sys.exit(2)
    with open(captions_path) as f:
        captions = json.load(f)
    captioned_dir = os.path.join(work, "captioned_decks")
    audit_dir = os.path.join(work, "audit")

    print(f"{'Deck':<34}{'Pics':>5}{'Cap':>4}{'Skp':>4}"
          f"{'T':>3}{'B':>3}{'C':>3}{'E':>3}{'O':>3}{'P':>3}{'Miss':>5}{'OK':>5}")
    print("=" * 79)

    grand = {"decks": 0, "T": 0, "B": 0, "C": 0, "E": 0, "O": 0, "P": 0,
             "text_miss": 0, "skips": 0, "unhashable": 0}
    for path in sorted(glob.glob(os.path.join(captioned_dir, "*_captioned.pptx"))):
        deck = os.path.basename(path).replace("_captioned.pptx", "")
        known_skip = load_known_skips(audit_dir, deck)
        prs = Presentation(path)
        H = prs.slide_height
        # v0.2.5: repeated-background hashes (same image on >= BG_REPEAT_THRESHOLD
        # slides) are NOT picture obstacles — apply lets captions sit over a
        # full-bleed background, so verify must exclude them from the Pattern C/P
        # "inside a picture" set or it false-flags every such caption. Mirrors the
        # bg_hashes pre-pass in apply_captions.apply_to_deck.
        _bg_slidecount = {}
        for _si, _sl in enumerate(prs.slides):
            for _p in iter_slide_pics(_sl):
                if _p["rid"] is None:
                    continue
                try:
                    _hh = hashlib.sha256(resolve_blob(_sl, _p["rid"])).hexdigest()[:12]
                except Exception:
                    continue
                _bg_slidecount.setdefault(_hh, set()).add(_si)
        bg_hashes_v = {h for h, s in _bg_slidecount.items()
                       if len(s) >= BG_REPEAT_THRESHOLD}
        total_pics = total_match = 0
        t_fail = b_fail = c_fail = e_fail = o_fail = p_fail = skip_known = unhashable = 0
        not_captioned = 0   # decorative or absent-from-captions.json — not a miss
        for slide in prs.slides:
            flim = slide_footer_top(slide, H, exclude_caption_shapes=True)
            text_rects = slide_text_obstacle_rects(slide)  # title + body + text boxes
            textboxes = []
            for sh in slide.shapes:
                if sh.has_text_frame:
                    t = (sh.text_frame.text or "").strip()
                    if t:
                        textboxes.append(t)
            # Picture geometry list for Pattern C check on this slide
            slide_pic_rects = []
            for pic in iter_slide_pics(slide):
                if pic["rid"] is None:
                    continue
                total_pics += 1
                try:
                    blob = resolve_blob(slide, pic["rid"])
                    ext = guess_ext(slide.part.related_part(pic["rid"]))
                    h = hashlib.sha256(blob).hexdigest()[:12]
                except Exception:
                    unhashable += 1   # surfaced in TOTAL; was a silent drop pre-v0.2.2
                    continue
                ox, oy, cx, cy = (pic.get("off_x"), pic.get("off_y"),
                                  pic.get("ext_cx"), pic.get("ext_cy"))
                if None not in (ox, oy, cx, cy) and h not in bg_hashes_v:
                    slide_pic_rects.append((int(ox), int(oy), int(cx), int(cy)))
                key = f"{deck}/{h}.{ext}"
                expected = captions.get(key)
                if not expected or expected.strip().lower() in ("[decorative]", "decorative"):
                    not_captioned += 1   # intentionally uncaptioned — exclude from miss
                    continue
                if h in known_skip:
                    skip_known += 1
                    continue  # intentionally not placed
                if expected in textboxes:
                    total_match += 1

            # Collect caption shapes once, then run all coordinate checks.
            caps = list(caption_shapes(slide))
            for nm, (cl, ct, cw, ch), txt, fpt in caps:
                is_sa_icon = nm.startswith(SA_ICON_PREFIX)
                is_band = nm.startswith(BAND_PREFIX)
                # Pattern O (v0.2.4): a single word wider than the box -> overflow.
                # Use the caption's real font size (fpt) so a down-scaled icon
                # caption is measured at its actual size, not a fixed 8pt.
                if caption_overflows(txt, cw, is_icon=is_sa_icon, font_pt=fpt):
                    o_fail += 1
                # Pattern P (v0.2.4): a band caption burying a small/structural picture.
                if is_band:
                    _cap_ltrb = (cl, ct, cl + cw, ct + ch)
                    _best = None; _bov = 0
                    for pr in slide_pic_rects:
                        _ov = rect_intersect_area((cl, ct, cw, ch), pr)
                        if _ov > _bov:
                            _bov = _ov; _best = pr
                    if _best is not None and _bov > 0:
                        _pic_ltrb = (_best[0], _best[1], _best[0] + _best[2], _best[1] + _best[3])
                        if band_covers_structural_picture(_cap_ltrb, _pic_ltrb):
                            p_fail += 1
                # Pattern T: caption overlaps ANY text frame (title/body/text box)
                # by >5% of the caption's area. Applies to every caption, INCLUDING
                # band captions (a band may sit in a picture but must never cover text).
                _cap_area = max(1, cw * ch)
                for tr in text_rects:
                    ix = min(cl + cw, tr[2]) - max(cl, tr[0])
                    iy = min(ct + ch, tr[3]) - max(ct, tr[1])
                    if ix > 0 and iy > 0 and ix * iy > 0.05 * _cap_area:
                        t_fail += 1
                        break
                # Pattern B: caption bottom too close to footer.
                if (ct + ch) > flim - FOOTER_CLEARANCE_EMU:
                    b_fail += 1
                # Pattern C: caption mostly inside a picture. EXEMPT: sa_icons
                # (inside their diagram) and band captions (deliberately in the
                # own-picture bottom strip — there was no text-clear external slot).
                if not (is_sa_icon or is_band):
                    for pr in slide_pic_rects:
                        ov = rect_intersect_area((cl, ct, cw, ch), pr)
                        if cw * ch > 0 and ov / (cw * ch) > 0.50:
                            c_fail += 1
                            break
            # Pattern E: pairwise caption-caption overlap (>5% of either box).
            for i in range(len(caps)):
                for j in range(i + 1, len(caps)):
                    a = caps[i][1]; b = caps[j][1]   # each is (l, t, w, h)
                    ov = rect_intersect_area(a, b)   # rect_intersect_area wants (l,t,w,h)
                    aa = a[2] * a[3]; bb = b[2] * b[3]
                    if ov > 0 and ((aa and ov / aa > 0.05) or (bb and ov / bb > 0.05)):
                        e_fail += 1

        text_miss = max(0, total_pics - total_match - skip_known
                        - unhashable - not_captioned)
        n_caps = sum(1 for _ in
                     (c for s in prs.slides for c in caption_shapes(s)))
        ok = "✓" if (t_fail == b_fail == c_fail == e_fail == o_fail == p_fail == 0
                     and text_miss == 0) else "✗"
        print(f"{deck[:33]:<34}{total_pics:>5}{n_caps:>4}{skip_known:>4}"
              f"{t_fail:>3}{b_fail:>3}{c_fail:>3}{e_fail:>3}{o_fail:>3}{p_fail:>3}"
              f"{text_miss:>5}{ok:>5}")
        grand["decks"] += 1
        for k, v in (("T", t_fail), ("B", b_fail), ("C", c_fail),
                     ("O", o_fail), ("P", p_fail),
                     ("E", e_fail), ("text_miss", text_miss),
                     ("skips", skip_known), ("unhashable", unhashable)):
            grand[k] += v

    print("=" * 79)
    print(f"TOTAL: {grand['decks']} decks | "
          f"T(text-overlap)={grand['T']} B(footer)={grand['B']} O(overflow)={grand['O']} P(covers-pic)={grand['P']} "
          f"C(in-pic)={grand['C']} E(cap-cap)={grand['E']} | "
          f"text-miss={grand['text_miss']} | known-skips={grand['skips']} | "
          f"unhashable={grand['unhashable']}")
    placement_fail = any(grand[k] for k in ("T", "B", "C", "E", "O", "P"))
    if grand["text_miss"] > 0:
        sys.exit(5)
    sys.exit(4 if placement_fail else 0)


if __name__ == "__main__":
    main()
