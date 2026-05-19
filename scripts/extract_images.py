"""Extract unique images from .pptx (or all .pptx in a folder) into a work dir.

Usage:
  python3 extract_images.py <input.pptx-or-folder> <work_dir> [--context "Biology 101 Lecture 3"]

Produces:
  <work_dir>/manifest.json
  <work_dir>/images/<deck_name>/<sha256-hash>.<ext>

Features (improvements from Gemini review v1/v2):
  - Recursive traversal into GROUP shapes (catches nested pictures consultants love to group)
  - Skips hidden shapes (cNvPr hidden="1")
  - Per-image progress prints
  - Try/except wrapper per deck (one corrupt/password-protected file does not halt the batch)
  - Captures slide title + body text context for caption quality
  - Optional --context flag for deck-level overarching context (e.g., course name)
"""
import sys, os, json, hashlib, glob, argparse, traceback
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from _oxml_pics import iter_slide_pics, resolve_blob, guess_ext, NS

# NOTE (adversarial-review fix): picture enumeration moved off the
# python-pptx `shape_type == PICTURE` walk and onto a strict raw-OOXML
# slide-spTree `.//p:pic` enumeration (see _oxml_pics.py). The old walk typed
# placeholder-hosted and OLE-fallback `<p:pic>` as PLACEHOLDER and silently
# dropped them (505 ~12%, 507 88%, 534 97% coverage). Hashes stay byte-identical
# (slide.part.related_part(rId).blob == old pic.image.blob, proven end-to-end).


def extract_text_context(slide):
    """Pull all visible text on a slide for caption context (title + body)."""
    bits = []
    for sh in slide.shapes:
        try:
            if sh.has_text_frame:
                t = (sh.text_frame.text or '').strip()
                if t:
                    bits.append(t)
        except Exception:
            continue
    return ' | '.join(bits)[:400]  # cap


def process_deck(deck_path, images_root, deck_context=''):
    deck = os.path.splitext(os.path.basename(deck_path))[0]
    deck_dir = os.path.join(images_root, deck)
    os.makedirs(deck_dir, exist_ok=True)
    try:
        prs = Presentation(deck_path)
    except Exception as e:
        print(f"  ERROR opening {deck}: {type(e).__name__}: {str(e)[:80]}")
        return {
            'deck': deck, 'deck_path': os.path.abspath(deck_path),
            'error': f'{type(e).__name__}: {e}', 'n_slides': 0, 'n_pictures': 0,
            'n_unique_images': 0, 'pictures': [], 'deck_context': deck_context,
        }

    seen = {}
    pictures = []
    n_groups_seen = 0
    n_linked_skipped = 0  # <a:blip> with r:link but no r:embed (external pic)

    for s_idx, slide in enumerate(prs.slides, 1):
        slide_text = extract_text_context(slide)
        # Count top-level groups for diagnostics
        for sh in slide.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                n_groups_seen += 1

        for pic in iter_slide_pics(slide):
            depth = pic['depth']
            if pic['rid'] is None:
                # No r:embed. Linked/external picture (or malformed blip): cannot
                # hash bytes -> SKIP + structured log (parity with apply/verify).
                n_linked_skipped += 1
                print(f"  SKIP slide {s_idx} pic id={pic['pic_id']} "
                      f"name={pic['pic_name']!r}: "
                      f"{'linked (r:link, no r:embed)' if pic['linked'] else 'no r:embed on a:blip'}")
                continue
            try:
                blob = resolve_blob(slide, pic['rid'])
                ext = guess_ext(slide.part.related_part(pic['rid']))
            except Exception as e:
                print(f"  SKIP slide {s_idx} pic id={pic['pic_id']}: "
                      f"blob resolve failed ({type(e).__name__})")
                continue
            h = hashlib.sha256(blob).hexdigest()[:12]
            pic_id = pic['pic_id']; pic_name = pic['pic_name']
            old_descr = pic['old_descr']
            fname = f"{h}.{ext}"
            if h not in seen:
                with open(os.path.join(deck_dir, fname), 'wb') as f:
                    f.write(blob)
                seen[h] = fname
            pictures.append({
                'slide': s_idx, 'pic_id': pic_id, 'name': pic_name,
                'image_hash': h, 'image_file': fname, 'ext': ext,
                'bytes': len(blob), 'in_group_depth': depth,
                'slide_text_context': slide_text, 'old_descr': old_descr,
            })

    return {
        'deck': deck, 'deck_path': os.path.abspath(deck_path),
        'n_slides': len(prs.slides), 'n_pictures': len(pictures),
        'n_unique_images': len(seen), 'n_top_level_groups_seen': n_groups_seen,
        'n_linked_skipped': n_linked_skipped,
        'pictures': pictures, 'deck_context': deck_context,
    }


def main():
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument('input', help='Single .pptx or a directory of .pptx files')
    ap.add_argument('work_dir', help='Output working directory')
    ap.add_argument('--context', default='', help='Deck-level overarching context (e.g., "Biology 101, Lecture 3: Mitosis")')
    args = ap.parse_args()

    inp = os.path.abspath(args.input)
    work = os.path.abspath(args.work_dir)
    os.makedirs(work, exist_ok=True)
    images_root = os.path.join(work, 'images')
    os.makedirs(images_root, exist_ok=True)

    if os.path.isfile(inp):
        decks = [inp]
    elif os.path.isdir(inp):
        decks = sorted(glob.glob(os.path.join(inp, '*.pptx')))
        # Filter out our own output to avoid recursion
        decks = [d for d in decks if not d.endswith('_captioned.pptx')]
    else:
        print(f"ERROR: {inp} not found"); sys.exit(2)

    if not decks:
        print("No .pptx files found."); sys.exit(2)

    manifest = {}
    for i, d in enumerate(decks, 1):
        print(f"[{i}/{len(decks)}] {os.path.basename(d)}")
        info = process_deck(d, images_root, args.context)
        manifest[info['deck']] = info
        if 'error' in info:
            print(f"           SKIPPED (error captured in manifest)")
        else:
            groups = info.get('n_top_level_groups_seen', 0)
            extra = f"  (+{groups} group shapes traversed)" if groups else ""
            print(f"           {info['n_slides']:>3} slides  {info['n_pictures']:>3} pics  {info['n_unique_images']:>3} unique{extra}")

    out = os.path.join(work, 'manifest.json')
    with open(out, 'w') as f:
        json.dump(manifest, f, indent=2)

    n_ok = sum(1 for v in manifest.values() if 'error' not in v)
    n_err = sum(1 for v in manifest.values() if 'error' in v)
    print(f"\nManifest: {out}  ({n_ok} ok, {n_err} errored)")
    if args.context:
        print(f"Deck context: '{args.context}' applied to all decks")
    print(f"Images:   {images_root}/")
    print(f"\nNext: Claude reads each image and writes captions.json to {work}/captions.json")


if __name__ == '__main__':
    main()
