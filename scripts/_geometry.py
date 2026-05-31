"""Shared geometry helpers for the captioner placement and audit pipelines.

Imported by apply_captions.py, audit_caption_placement.py, and verify.py so
all three agree on slide-band thresholds, footer clearance, and obstacle-rect
semantics. Keeping these constants/functions in one place prevents drift —
which previously caused audit-vs-apply geometric mismatches.

v0.2.2 adds true 2D box-intersection obstacle checking (`clear_all_obstacles_2d`
+ `slide_body_obstacle_rects`): a caption is only "blocked" by an obstacle when
it overlaps in BOTH axes. This fixes two v0.2.1 defects: (1) vertical-only
checks falsely blocked captions that were horizontally clear of a narrow
title/body placeholder (forcing bad fallbacks), and (2) caption-caption
avoidance was silently disabled because placed-caption rects (l,t,w,h) were
misread as (L,T,R,B). The legacy vertical-only helpers (`_clear_all_obstacles`,
`_voverlap`, `slide_body_obstacle_bands`) are retained for back-compat.
"""
import os
from pptx.enum.shapes import PP_PLACEHOLDER

# ---------------------------------------------------------------------------
# Geometry constants — single source of truth for all three pipelines.
# ---------------------------------------------------------------------------
MIN_CAPTION_HEIGHT  = 250_000   # ~0.27 in — absolute floor for a readable caption box.
FOOTER_CLEARANCE_EMU = 91_440   # ~0.10 in — visual gap between caption bottom and footer band.
MIN_CAPTION_WIDTH_EMU = 1_828_800  # ~2.00 in — minimum width so a 30-50 char caption fits in 1 line.
EMU_PER_CHAR_DEFAULT  = 114_000  # ~0.125 in/char — conservative estimate for 10pt italic Calibri.
                                 # See "Fix-D heuristic" in SKILL.md "Known limits."
LOW_BAND_FRACTION   = 0.86       # bottom 14% of slide = "footer band" search zone.

# Placeholder type sets — kept here so apply/audit/verify agree on what counts
# as a "title" obstacle vs a "body" obstacle vs a "footer" land-mine.
TITLE_TYPES  = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
BODY_TYPES   = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.SUBTITLE,
                PP_PLACEHOLDER.OBJECT}  # CONTENT alias = OBJECT in python-pptx
FOOTER_TYPES = {PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.DATE,
                PP_PLACEHOLDER.SLIDE_NUMBER}

# captioner_capband_ = a caption deliberately placed in the bottom strip of its
# OWN picture (v0.2.3 fallback) — better than covering text or skipping. The
# audit / verify in-picture checks exempt this prefix (the overlap is intended).
CAPTION_NAME_PREFIXES = ('captioner_caption_', 'captioner_sa_icon_', 'captioner_capband_')


# ---------------------------------------------------------------------------
# Footer band: the y-coordinate below which captions must not extend.
# ---------------------------------------------------------------------------
def slide_footer_top(slide, slide_h, exclude_caption_shapes=True):
    """Return the topmost y (EMU) of the reserved bottom band. Considers
    FOOTER/DATE/SLIDE_NUMBER placeholders on slide/layout/master AND any
    low-band text shape (catches branded master footers like "Pearson").
    A footer-typed placeholder PARKED AT THE TOP of a master (top < band)
    is NOT a bottom footer — ignore it (else the band collapses to the
    slide top). When called on a captioned deck, caption shapes are
    excluded so they don't pollute their own footer line.
    """
    limit = slide_h - 274_320  # ~0.30 in universal bottom safety margin
    band  = int(slide_h * LOW_BAND_FRACTION)

    sources = [slide]
    try:
        sources.append(slide.slide_layout)
        sources.append(slide.slide_layout.slide_master)
    except Exception:
        pass

    for src in sources:
        try:
            shapes = list(src.shapes)
        except Exception:
            continue
        for sh in shapes:
            try:
                if exclude_caption_shapes:
                    nm = sh.name or ''
                    if any(nm.startswith(p) for p in CAPTION_NAME_PREFIXES):
                        continue
                top = sh.top
                if top is None or top < band:
                    continue
                is_foot = False
                try:
                    pf = sh.placeholder_format
                    if pf is not None and pf.type in FOOTER_TYPES:
                        is_foot = True
                except Exception:
                    pass
                has_txt = bool(getattr(sh, 'has_text_frame', False)
                               and sh.text_frame.text.strip())
                if is_foot or has_txt:
                    limit = min(limit, int(top))
            except Exception:
                continue
    return max(0, limit)


# ---------------------------------------------------------------------------
# Title obstacle: caption must not visually overlap a slide title.
# ---------------------------------------------------------------------------
def slide_title_rect(slide):
    """Bounding rect (L, T, R, B) EMU enclosing every TITLE/CENTER_TITLE
    placeholder with visible text, wherever positioned. None if no titled text."""
    L = T = R = B = None
    for sh in slide.shapes:
        if (sh.name or '').startswith(CAPTION_NAME_PREFIXES):
            continue
        try:
            pf = sh.placeholder_format
            if pf is None or pf.type not in TITLE_TYPES:
                continue
            if not (getattr(sh, 'has_text_frame', False)
                    and sh.text_frame.text.strip()):
                continue
            if None in (sh.left, sh.top, sh.width, sh.height):
                continue
            l, t = int(sh.left), int(sh.top)
            r, b = l + int(sh.width), t + int(sh.height)
            L = l if L is None else min(L, l)
            T = t if T is None else min(T, t)
            R = r if R is None else max(R, r)
            B = b if B is None else max(B, b)
        except Exception:
            continue
    return None if L is None else (L, T, R, B)


# ---------------------------------------------------------------------------
# Body obstacles: SUBTITLE/BODY/OBJECT placeholders with non-empty text.
# Returns vertical-band tuples (y_top, y_bottom). Vertical-only is the
# v0.2.1 contract — v0.3.0 may upgrade to 2D rects.
# ---------------------------------------------------------------------------
def slide_body_obstacle_bands(slide):
    """List of (y_top, y_bottom) for every body-type placeholder with visible
    text on this slide. Used by the placement obstacle filter. Vertical-only
    semantics (consistent with _voverlap)."""
    bands = []
    for sh in slide.shapes:
        if (sh.name or '').startswith(CAPTION_NAME_PREFIXES):
            continue
        try:
            pf = sh.placeholder_format
            if pf is None or pf.type not in BODY_TYPES:
                continue
            if not (getattr(sh, 'has_text_frame', False)
                    and sh.text_frame.text.strip()):
                continue
            if sh.top is None or sh.height is None:
                continue
            bands.append((int(sh.top), int(sh.top) + int(sh.height)))
        except Exception:
            continue
    return bands


def slide_body_obstacle_rects(slide):
    """List of (L, T, R, B) EMU rects for every body-type placeholder with
    visible text on this slide. The 2D successor to slide_body_obstacle_bands
    — carries horizontal extent so clear_all_obstacles_2d can tell a narrow
    centered subtitle apart from a right-side caption that doesn't touch it.
    """
    rects = []
    for sh in slide.shapes:
        if (sh.name or '').startswith(CAPTION_NAME_PREFIXES):
            continue
        try:
            pf = sh.placeholder_format
            if pf is None or pf.type not in BODY_TYPES:
                continue
            if not (getattr(sh, 'has_text_frame', False)
                    and sh.text_frame.text.strip()):
                continue
            if None in (sh.left, sh.top, sh.width, sh.height):
                continue
            l, t = int(sh.left), int(sh.top)
            rects.append((l, t, l + int(sh.width), t + int(sh.height)))
        except Exception:
            continue
    return rects


def _visible_text_rect(sh):
    """(L, T, R, B) approximating the VISIBLE text region of a text shape, not its
    full bounding box. A title placeholder is often a tall box with one top-anchored
    line of text — using the full bbox makes a caption sitting below the title text
    (but inside the empty lower bbox) look like an overlap. Estimate the text height
    from paragraph/character counts and anchor it (top/middle/bottom) within the box.
    Conservative: never smaller than one line, never larger than the box."""
    l, t, w, hgt = int(sh.left), int(sh.top), int(sh.width), int(sh.height)
    try:
        tf = sh.text_frame
        # Font-aware line counting. A fixed 10pt char width badly under-counts
        # lines for large fonts (a 37pt title wraps in ~1/4 the chars), and
        # paragraphs carry explicit line breaks (\n, and \x0b/\v for a soft
        # Shift+Enter break) that must each start a new visual line. Under-counting
        # lines shrinks the obstacle so a caption lands on the title's 2nd line.
        sizes = [r.font.size.pt for para in tf.paragraphs for r in para.runs
                 if r.font.size is not None]
        font_pt = max(sizes) if sizes else 18.0
        char_w = max(1.0, font_pt * _EM_FRACTION * _EMU_PER_PT)
        cpl = max(4, int(w / char_w))
        lines = 0
        for para in tf.paragraphs:
            txt = (para.text or '')
            if not txt.strip():
                continue
            for seg in txt.replace('\x0b', '\n').replace('\v', '\n').split('\n'):
                lines += max(1, (len(seg) + cpl - 1) // cpl) if seg.strip() else 1
        lines = max(1, lines)
        line_emu = int(font_pt * 1.3 * _EMU_PER_PT)  # font-scaled line height
        vis_h = min(hgt, lines * line_emu + 40_000)  # + small padding
        anchor = None
        try:
            anchor = tf.vertical_anchor  # MSO_ANCHOR or None (None ~ top)
        except Exception:
            anchor = None
        a = str(anchor) if anchor is not None else 'TOP'
        if 'BOTTOM' in a:
            vt = t + hgt - vis_h
        elif 'MIDDLE' in a or 'CENTER' in a:
            vt = t + (hgt - vis_h) // 2
        else:                       # TOP / unknown / inherited
            vt = t
        return (l, vt, l + w, vt + vis_h)
    except Exception:
        return (l, t, l + w, t + hgt)


def slide_text_obstacle_rects(slide):
    """List of (L, T, R, B) EMU rects for the VISIBLE text region of EVERY shape on
    the slide that carries text — title, body/subtitle, AND ordinary text boxes /
    auto-shapes (AUTO_SHAPE, TEXT_BOX) that are NOT placeholders. v0.2.3 closed the
    gap where captions covered text in plain text boxes; v0.2.3.1 narrows each
    obstacle to its estimated visible-text region (anchor-aware) so a caption below
    a top-anchored title's text isn't falsely blocked by the title's tall empty box.
    Excludes captioner's own shapes and footer/date/slide-number placeholders.
    """
    rects = []
    for sh in slide.shapes:
        try:
            nm = sh.name or ''
            if nm.startswith(CAPTION_NAME_PREFIXES):
                continue
            try:
                pf = sh.placeholder_format
                if pf is not None and pf.type in FOOTER_TYPES:
                    continue
            except Exception:
                pass
            if not (getattr(sh, 'has_text_frame', False)
                    and sh.text_frame.text.strip()):
                continue
            if None in (sh.left, sh.top, sh.width, sh.height):
                continue
            rects.append(_visible_text_rect(sh))
        except Exception:
            continue
    return rects


def estimate_caption_height(text, box_width_emu, line_height_emu=160_000,
                            margin_emu=40_000, nominal_emu=MIN_CAPTION_HEIGHT):
    """Estimate the rendered height of an auto-sizing caption box so placement and
    overlap math use the height the box will ACTUALLY occupy (a 2-line caption in a
    1-line box otherwise grows down onto whatever is below it). Conservative: floors
    at `nominal_emu`, assumes ~EMU_PER_CHAR_DEFAULT per character for wrap."""
    if not text:
        return nominal_emu
    cpl = max(8, int(box_width_emu // EMU_PER_CHAR_DEFAULT))
    lines = max(1, (len(text) + cpl - 1) // cpl)
    return max(nominal_emu, lines * line_height_emu + margin_emu)


# ---------------------------------------------------------------------------
# Vertical overlap predicates.
# ---------------------------------------------------------------------------
def _voverlap_band(c_top, c_h, y_top, y_bottom, frac=0.15):
    """True if caption interval [c_top, c_top+c_h] vertically overlaps the
    band [y_top, y_bottom] by more than `frac` of caption height."""
    if y_top is None or y_bottom is None:
        return False
    iy = max(0, min(c_top + c_h, y_bottom) - max(c_top, y_top))
    return iy > frac * c_h


def _voverlap(c_top, c_h, title_rect, frac=0.15):
    """Back-compat shim for the original signature. title_rect = (L, T, R, B)."""
    if title_rect is None:
        return False
    return _voverlap_band(c_top, c_h, title_rect[1], title_rect[3], frac)


def _clear_all_obstacles(c_top, c_h, obstacles, frac=0.15):
    """True if caption [c_top, c_top+c_h] clears EVERY obstacle. Obstacles
    accepted in either format:
      - 4-tuple (L, T, R, B): use T, B (vertical-only — v0.2.1 contract)
      - 2-tuple (y_top, y_bottom): use directly
    """
    for obs in obstacles:
        if obs is None:
            continue
        if len(obs) == 4:
            y_top, y_bottom = obs[1], obs[3]
        elif len(obs) == 2:
            y_top, y_bottom = obs
        else:
            continue
        if _voverlap_band(c_top, c_h, y_top, y_bottom, frac):
            return False
    return True


# ---------------------------------------------------------------------------
# 2D obstacle clearance (v0.2.2) — require overlap in BOTH axes to block.
# ---------------------------------------------------------------------------
def _to_ltrb(obs):
    """Normalize an obstacle into an (L, T, R, B) rect, or None if it can't be.
      - 4-tuple already (L, T, R, B)  -> as-is
      - 2-tuple (y_top, y_bottom)     -> full-width band (no horizontal info):
        return None so the caller can fall back to a vertical-only test rather
        than guess an x-extent. (All v0.2.2 apply-side obstacles are 4-tuples.)
    """
    if obs is None:
        return None
    if len(obs) == 4:
        return (obs[0], obs[1], obs[2], obs[3])
    return None


def _rects_block(cap, obs, vfrac=0.15):
    """True if caption rect `cap`=(L,T,R,B) is BLOCKED by obstacle `obs`=(L,T,R,B):
    they share horizontal extent (ix > 0) AND vertically overlap by more than
    `vfrac` of the caption's height. Both conditions required — a horizontally
    disjoint obstacle never blocks (the core v0.2.2 fix)."""
    cl, ct, cr, cb = cap
    ol, ot, orr, ob = obs
    ix = min(cr, orr) - max(cl, ol)
    iy = min(cb, ob) - max(ct, ot)
    c_h = max(1, cb - ct)
    return ix > 0 and iy > vfrac * c_h


def clear_all_obstacles_2d(cap_rect, obstacles, vfrac=0.15):
    """True if caption `cap_rect`=(L,T,R,B) clears EVERY obstacle in 2D.
    A 4-tuple obstacle is tested in 2D; a 2-tuple (y_top,y_bottom) obstacle,
    which lacks horizontal extent, falls back to the legacy vertical-only test
    (conservative — treats it as full-width)."""
    cl, ct, cr, cb = cap_rect
    c_h = max(1, cb - ct)
    for obs in obstacles:
        r = _to_ltrb(obs)
        if r is None:
            # No horizontal info — vertical-only fallback (full-width band).
            if obs is not None and len(obs) == 2:
                if _voverlap_band(ct, c_h, obs[0], obs[1], vfrac):
                    return False
            continue
        if _rects_block(cap_rect, r, vfrac):
            return False
    return True


# ---------------------------------------------------------------------------
# 2D rect intersection — used by audit for caption-caption overlap detection.
# ---------------------------------------------------------------------------
def rect_intersect_area(a, b):
    """Area of intersection of two (l, t, w, h) rects. 0 if disjoint."""
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0, min(ay + ah, by + bh) - max(ay, by))
    return ix * iy


# ---------------------------------------------------------------------------
# v0.2.4 — text-overflow + structural-picture-coverage detection.
# A single word cannot wrap, so a caption box narrower than its longest word
# overflows horizontally. Char width uses a conservative em-fraction LOWER bound
# (0.50 of font size) so the check never under-reports. Icon captions render at
# 8pt / 20000-EMU side margins; regular + band captions at 10pt / 50000.
# ---------------------------------------------------------------------------
_EMU_PER_PT = 12700.0
_EM_FRACTION = 0.50

def _char_w(is_icon, font_pt=None):
    # font_pt overrides the nominal size — used when an icon caption is scaled
    # DOWN to match a small icon (so the width math tracks the real glyph size).
    pt = font_pt if font_pt else (8.0 if is_icon else 10.0)
    return pt * _EM_FRACTION * _EMU_PER_PT  # 50800 @8pt, 63500 @10pt

def required_caption_width(text, is_icon=False, headroom=1.20, font_pt=None):
    """Box width (EMU) that comfortably fits the longest single word on one line.
    `_char_w` is a deliberate LOWER bound (0.50 em); real italic glyphs run wider,
    so SIZING adds `headroom` (default +20%) above the detection threshold — the
    box ends up comfortably wider than the word, not exactly at the edge. The
    `caption_overflows` gate keeps using the bare bound so it never false-flags.
    `font_pt` lets a down-scaled icon caption size its box to the smaller glyphs."""
    words = (text or "").split()
    if not words:
        return 0
    ml = 20000 if is_icon else 50000
    return int(max(len(w) for w in words) * _char_w(is_icon, font_pt) * headroom + 2 * ml)

def caption_overflows(text, box_width, is_icon=False, font_pt=None):
    """True if the longest word in `text` cannot fit on one line in `box_width`
    (horizontal overflow — the box must be widened or the caption shortened).
    `font_pt` must match the caption's actual rendered size so the gate agrees
    with how the box was sized in apply (else a down-scaled icon false-flags)."""
    words = (text or "").split()
    if not words:
        return False
    ml = 20000 if is_icon else 50000
    return max(len(w) for w in words) * _char_w(is_icon, font_pt) > (box_width - 2 * ml)

def band_covers_structural_picture(cap_ltrb, pic_ltrb):
    """True if a bottom-band caption AT (cap_ltrb) over picture (pic_ltrb) is
    AT-RISK of burying meaningful picture content — i.e. the picture is small,
    thin, or a tall-narrow structural column (numbered chevron / icon / banner),
    as opposed to a large photo where a thin bottom strip is harmless. Thresholds
    from the 49-deck corpus scan (2026-05-30)."""
    pw = pic_ltrb[2] - pic_ltrb[0]
    ph = pic_ltrb[3] - pic_ltrb[1]
    if pw <= 0 or ph <= 0:
        return False
    frac_h = (cap_ltrb[3] - cap_ltrb[1]) / ph
    inch = 914400.0
    min_edge = min(pw, ph)
    ar = pw / ph
    if frac_h > 0.40:
        return True
    if min_edge < 1.6 * inch and frac_h > 0.25:
        return True
    if ar < 0.85 and ph < 2.5 * inch and frac_h > 0.25:
        return True
    return False


# ---------------------------------------------------------------------------
# Picture coverage (for informational metadata in overlay-fullbleed rows).
# ---------------------------------------------------------------------------
def visible_coverage(p_left, p_top, p_width, p_height, slide_w, slide_h):
    """Fraction of slide area that the picture's on-slide-visible region
    covers. Clamps to slide bounds (a picture that overflows the slide
    edges doesn't get credit for the off-slide portion)."""
    vis_left   = max(0, p_left)
    vis_top    = max(0, p_top)
    vis_right  = min(slide_w, p_left + p_width)
    vis_bottom = min(slide_h, p_top + p_height)
    vis_w = max(0, vis_right - vis_left)
    vis_h = max(0, vis_bottom - vis_top)
    return (vis_w * vis_h) / (slide_w * slide_h)
