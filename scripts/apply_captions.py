"""Apply captions from captions.json to every deck listed in manifest.json.

(Side-effect imports: re used by SmartArt icon-name extractor below.)


Usage:
  python3 apply_captions.py <work_dir> [--dry-run] [--font-name FONT]
                            [--font-size PT] [--font-color HEX] [--italic BOOL]
                            [--gap-emu INT] [--height-emu INT]
                            [--update-existing] [--quiet]

Expects:
  <work_dir>/manifest.json    (from extract_images.py)
  <work_dir>/captions.json    (written by Claude after reading images)
      Format: {"<deck>/<hash>.<ext>": "<short caption>", ...}
      Special value "[decorative]" → skip (no visible text added, logged decorative)

Produces:
  <work_dir>/captioned_decks/<deck>_captioned.pptx
  <work_dir>/audit/<deck>_audit.csv

Improvements (from Gemini review v1/v2):
  - --dry-run: emits audit CSV but does NOT modify any .pptx
  - --font-name / --font-size / --font-color / --italic: caption styling
  - --gap-emu / --height-emu: layout offsets
  - Recursive traversal into GROUP shapes
  - Re-run protection: detects `_captioned` suffix on input + warns
  - Idempotency: caption shapes get name prefix "captioner_caption_<hash>" so re-runs can
    detect + with --update-existing flag, remove and re-add (else: append, leading to duplicates)
  - Per-deck progress prints
  - Per-deck try/except — corrupt/password-protected file does not halt batch
"""
import sys, os, json, shutil, csv, hashlib, argparse, re
from collections import Counter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Pt
from _oxml_pics import iter_slide_pics, resolve_blob, guess_ext
from _geometry import (
    slide_footer_top as _g_slide_footer_top,
    slide_text_obstacle_rects as _g_text_rects,
    clear_all_obstacles_2d, visible_coverage, estimate_caption_height,
    required_caption_width, caption_overflows, band_covers_structural_picture,
    FOOTER_CLEARANCE_EMU, MIN_CAPTION_WIDTH_EMU, EMU_PER_CHAR_DEFAULT,
    MIN_CAPTION_HEIGHT as _G_MIN_CAPTION_HEIGHT,
)

# PACE player_b fix 2026-05-18: pic enumeration switched to the strict
# raw-OOXML slide-spTree `.//p:pic` walk (see _oxml_pics.py), matching
# extract/verify exactly. Placement geometry reads the pic's own
# <a:off>/<a:ext> EMU; verified byte-identical to the old pic.left/.top/
# .width/.height for every deck pic in this corpus (incl. group-nested),
# so caption-card positions do not move. Hashes unchanged.
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

CAPTION_SHAPE_NAME_PREFIX = 'captioner_caption_'  # idempotency marker
CAPTION_BAND_NAME_PREFIX = 'captioner_capband_'   # caption deliberately in own-picture bottom strip
BAND_FONT_PT = 6   # in-picture band captions render SMALL — overlap less photo, fit tight spots
SMARTART_SHAPE_NAME_PREFIX = 'captioner_smartart_'  # idempotency marker for SmartArt captions
SMARTART_ICON_SHAPE_NAME_PREFIX = 'captioner_sa_icon_'  # idempotency marker for per-icon SmartArt captions
MIN_CAPTION_HEIGHT = 250000

NS_A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
NS_DGM_REL = '{http://schemas.openxmlformats.org/drawingml/2006/diagram}'
DIAGRAM_DATA_URI = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData'

# ---------------------------------------------------------------------------
# Spell-check (opt-in, --spellcheck). FLAG-ONLY: this never edits the .pptx and
# never auto-corrects a word. It emits suspected misspellings to a separate
# <deck>_spellcheck.csv for a human to review. A bundled, user-extensible
# whitelist (spellcheck_whitelist.txt) prevents flagging known domain terms,
# abbreviations, and proper nouns — so captioner does not "fix" non-issues.
# ---------------------------------------------------------------------------
WHITELIST_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'spellcheck_whitelist.txt')
# Proper-noun typos a generic dictionary won't suggest well — surfaced explicitly.
# Names here are web-verified canonical spellings (see SKILL.md "Name verification").
SPELL_KNOWN_BAD = {
    'humaoid': 'humanoid', 'humaoids': 'humanoids',
    'appronik': 'Apptronik',          # Apptronik (humanoid-robotics co.) — web-verified
    'geoffery': 'Geoffrey',           # Geoffrey Moore (Crossing the Chasm) — web-verified
    'clayten': 'Clayton',             # Clayton Christensen (HBS) — web-verified
    'siemans': 'Siemens',             # Siemens — web-verified
    'wolfrom': 'Wolfram',             # Wolfram Research — web-verified
}


def _load_whitelist():
    wl = set()
    try:
        with open(WHITELIST_PATH, encoding='utf-8') as f:
            for line in f:
                line = line.strip()
                if not line or line.startswith('#'):
                    continue
                for tok in line.split():
                    wl.add(tok.lower())
    except OSError:
        pass
    return wl


def init_spellcheck(enabled):
    """Return (SpellChecker|None, whitelist set, status str). Degrades gracefully
    if pyspellchecker is not installed (feature is optional)."""
    if not enabled:
        return None, set(), 'disabled'
    wl = _load_whitelist()
    try:
        from spellchecker import SpellChecker
    except ImportError:
        return None, wl, 'unavailable'  # documented optional dependency
    return SpellChecker(distance=1), wl, 'ready'


def spell_scan(text, source, s_idx, sp, wl, seen, rows):
    """Append suspected-misspelling rows for `text`. FLAG-ONLY — no mutation.
    `source` is 'caption' or 'slide-text'. `seen` dedupes (source, slide, word)."""
    if not text:
        return
    low_line = text.lower()
    if 'http' in low_line or 'www.' in low_line or '://' in low_line \
            or any(d in low_line for d in ('.com', '.org', '.ai', '.io',
                                           '.gov', '.edu', '.net', '.co/')):
        return  # URL / citation / product-domain line: not prose
    for tok in re.findall(r"[A-Za-z][A-Za-z'\-]+", text):
        low = tok.lower().strip("'-")
        if low in SPELL_KNOWN_BAD:
            sug, known = SPELL_KNOWN_BAD[low], True
        else:
            if len(low) < 4 or tok.isupper():
                continue
            if any(ch.isdigit() for ch in tok) or "'" in tok or '-' in tok:
                continue
            # Plural-of-acronym (NPVs, IRRs, CEOs, KPIs, MNCs, SVMs, CNNs, GPTs):
            # the singular form is an all-caps acronym -> not a misspelling.
            if low.endswith('s') and len(tok) >= 3 and tok[:-1].isupper():
                continue
            if low in wl or sp is None:
                continue
            if not sp.unknown([low]):
                continue
            sug = sp.correction(low)
            if not sug or sug == low:
                continue
            known = False
        # Likely proper noun (Capitalized, not all-caps) -> the agent must
        # web-verify the canonical spelling before presenting it as a fix.
        verify_name = bool(tok[:1].isupper() and not tok.isupper())
        key = (source, s_idx, low)
        if key in seen:
            continue
        seen.add(key)
        rows.append({
            'slide': s_idx, 'source': source, 'term': tok,
            'suggestion': sug, 'known_bad': known,
            'verify_name': verify_name,
            'context': text.strip()[:110].replace('\n', ' '),
        })


_QC_PLACEHOLDER = ("click to add", "lorem ipsum", "[gap", "tbd", "xxx",
                   "placeholder text", "insert text here", "your text here")
_QC_DOUBLE = re.compile(r"\b(\w{3,})\s+\1\b", re.IGNORECASE)
_QC_DATE = re.compile(r"\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)"
                      r"[a-z]*\.?\s+\d{1,2}(?:,?\s*\d{4})?\b")


def qc_scan(text, source, s_idx, seen, rows):
    """READ-ONLY date / doubled-word / leftover-template scan. FLAG-ONLY — never
    edits the .pptx. Generic (no course-specific year/code rules). `seen` dedupes.
    Date strings are tagged 'date-review' (informational, not asserted wrong)."""
    if not text:
        return
    low = text.lower()
    for ph in _QC_PLACEHOLDER:
        if ph in low:
            k = ('placeholder', s_idx, ph)
            if k not in seen:
                seen.add(k)
                rows.append({'slide': s_idx, 'source': source, 'kind': 'leftover-template',
                             'detail': ph, 'context': text.strip()[:110].replace('\n', ' ')})
    for m in set(w.lower() for w in _QC_DOUBLE.findall(text)):
        if m in ('that', 'had', 'is', 'the'):  # common legitimate repeats
            continue
        k = ('double', s_idx, m)
        if k not in seen:
            seen.add(k)
            rows.append({'slide': s_idx, 'source': source, 'kind': 'doubled-word',
                         'detail': f"{m} {m}", 'context': text.strip()[:110].replace('\n', ' ')})
    for m in _QC_DATE.findall(text):
        k = ('date', s_idx, m.strip())
        if k not in seen:
            seen.add(k)
            rows.append({'slide': s_idx, 'source': source, 'kind': 'date-review',
                         'detail': m.strip(), 'context': text.strip()[:110].replace('\n', ' ')})


def resolve_ph_geometry(slide, ph_idx):
    """A picture content-placeholder's geometry is inherited from the layout/
    master; the <p:pic> itself often has NO <a:xfrm>. python-pptx resolves the
    inheritance, so match the slide placeholder by idx and read its effective
    box. Returns (left, top, width, height) EMU or None. (The 2026-05-18
    off-page-caption bug: placeholder pics fell back to (0, 50000).)"""
    try:
        for ph in slide.placeholders:
            pf = ph.placeholder_format
            if pf is not None and pf.idx == ph_idx:
                if None not in (ph.left, ph.top, ph.width, ph.height):
                    return int(ph.left), int(ph.top), int(ph.width), int(ph.height)
    except Exception:
        pass
    return None


def slide_footer_top(slide, slide_h):
    """Y (EMU) below which captions must NOT extend — the top of the reserved
    bottom band. Considers FOOTER/DATE/SLIDE_NUMBER placeholders on the slide,
    its layout, and master, plus any low-band text shape (catches branded
    footers like a 'Pearson' credit). Always reserves at least the bottom
    ~0.3in so captions never sit in the extreme bottom strip."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    foot_types = {PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.DATE,
                  PP_PLACEHOLDER.SLIDE_NUMBER}
    limit = slide_h - 274320  # ~0.30in universal bottom safety margin
    band = int(slide_h * 0.86)  # "low band" = bottom ~14%
    sources = [slide]
    try:
        sources.append(slide.slide_layout)
        sources.append(slide.slide_layout.slide_master)
    except Exception:
        pass
    for src in sources:
        try:
            shapes = list(src.placeholders) + [s for s in src.shapes
                                               if s not in src.placeholders]
        except Exception:
            try:
                shapes = list(src.shapes)
            except Exception:
                continue
        for sh in shapes:
            try:
                top = sh.top
                if top is None:
                    continue
                is_foot = False
                try:
                    pf = sh.placeholder_format
                    if pf is not None and pf.type in foot_types:
                        is_foot = True
                except Exception:
                    pass
                has_txt = bool(getattr(sh, 'has_text_frame', False)
                               and sh.text_frame.text.strip())
                # A footer is in the BOTTOM band. A footer/date/slide-number
                # placeholder parked at the TOP of a master (top < band) is
                # NOT a bottom footer — ignore it, else the reserved band
                # collapses to the slide top (the 2026-05-18 limit=113072 bug).
                if top >= band and (is_foot or has_txt):
                    limit = min(limit, int(top))
            except Exception:
                continue
    return max(0, limit)


def slide_title_box(slide):
    """Bounding box (l, t, r, b) EMU enclosing every TITLE/CENTER_TITLE
    placeholder that has visible text, wherever it sits on the slide (some
    layouts park the title low, under an image). None if no titled text.
    Captions must avoid this rect when a clear slot exists."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    title_types = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
    L = T = R = B = None
    for sh in slide.shapes:
        if (sh.name or '').startswith('captioner_caption_'):
            continue
        try:
            pf = sh.placeholder_format
            if pf is None or pf.type not in title_types:
                continue
            if not (getattr(sh, 'has_text_frame', False)
                    and sh.text_frame.text.strip()):
                continue
            if None in (sh.left, sh.top, sh.width, sh.height):
                continue
            l, t = int(sh.left), int(sh.top)
            r, b = l + int(sh.width), t + int(sh.height)
            L = l if L is None else min(L, l)
            T = t if T is None else min(T, t)
            R = r if R is None else max(R, r)
            B = b if B is None else max(B, b)
        except Exception:
            continue
    return None if L is None else (L, T, R, B)


def _voverlap(c_top, c_h, title):
    """True if a caption [c_top, c_top+c_h] vertically overlaps the title rect
    by a meaningful amount (>15% of caption height)."""
    if title is None:
        return False
    _, tT, _, tB = title
    iy = max(0, min(c_top + c_h, tB) - max(c_top, tT))
    return iy > 0.15 * c_h


def iter_slide_body_text(shapes):
    """Yield instructor-authored text (text frames + tables), recursing groups,
    skipping captioner's own added shapes."""
    for sh in shapes:
        try:
            name = sh.name or ''
        except Exception:
            name = ''
        if name.startswith((CAPTION_SHAPE_NAME_PREFIX, SMARTART_SHAPE_NAME_PREFIX,
                            SMARTART_ICON_SHAPE_NAME_PREFIX)):
            continue
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from iter_slide_body_text(sh.shapes)
                continue
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            yield cell.text
                continue
        except Exception:
            pass
        if getattr(sh, 'has_text_frame', False):
            t = sh.text_frame.text
            if t and t.strip():
                yield t


def iter_smartart_frames(slide):
    """Yield graphicFrame shapes that contain SmartArt (diagram) content."""
    for sh in slide.shapes:
        try:
            if not sh._element.tag.endswith('}graphicFrame'):
                continue
            for child in sh._element.iter():
                ctag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if ctag == 'graphicData' and 'diagram' in (child.get('uri') or ''):
                    yield sh
                    break
        except Exception:
            continue


def _camel_to_words(s):
    """Icons_VideoCamera -> 'video camera'. CheckMark -> 'check mark'."""
    s = s.replace('Icons_', '').replace('Icon_', '')
    # Insert space before each capital letter (except first), then lowercase
    out = re.sub(r'(?<!^)(?=[A-Z])', ' ', s).lower()
    return out.strip()


def extract_smartart_icons(slide_part, frame):
    """Extract icon names from SVG metadata embedded in the SmartArt's drawing part.

    Returns a deduped list of human-readable icon names in order of first appearance.
    Strategy: follow graphicFrame -> diagramDrawing rel -> read drawing*.xml's rels ->
    for each SVG image relationship, parse the SVG file and extract the `id` attribute
    (PowerPoint stores names like 'Icons_Checkmark' on the root <svg> element).
    """
    # Find rId of the diagramDrawing relationship
    drawing_rid = None
    for child in frame._element.iter():
        if child.tag.endswith('}relIds'):
            for k, v in child.attrib.items():
                # ms diagram-drawing rel is 'dgm' style — older r:dm references the data part,
                # but the drawing part is referenced via a 'http://schemas.microsoft.com/office/2007/relationships/diagramDrawing'
                # which lives under the diagramData's rels, NOT the slide's rels.
                if k.endswith('}dm') or k == 'dm':
                    drawing_rid = v
                    break
            if drawing_rid:
                break
    if not drawing_rid:
        return []
    try:
        # Slide -> diagramData part
        data_rel = slide_part.rels[drawing_rid]
        data_part = data_rel.target_part
        # Images may be linked directly from diagramData (common case) OR via a separate
        # diagramDrawing part (less common). Collect candidates from both.
        candidate_parts = [data_part]
        for rel in data_part.rels.values():
            if 'diagramDrawing' in rel.reltype:
                candidate_parts.append(rel.target_part)
        icons = []
        seen_targets = set()
        for cp in candidate_parts:
            for rel in cp.rels.values():
                if rel.reltype.endswith('/image'):
                    target = rel.target_part
                    pname = str(getattr(target, 'partname', '') or '').lower()
                    if not pname.endswith('.svg'):
                        continue
                    if pname in seen_targets:
                        continue
                    seen_targets.add(pname)
                    try:
                        blob = target.blob
                    except Exception:
                        continue
                    text = blob.decode('utf-8', errors='ignore')
                    m = re.search(r'<svg[^>]*\bid="([^"]+)"', text)
                    if not m:
                        continue
                    name = _camel_to_words(m.group(1))
                    if name and name not in icons:
                        icons.append(name)
        return icons
    except Exception:
        return []


def extract_smartart_icon_placements(slide_part, frame):
    """For each icon embedded in this SmartArt, return its name + slide-absolute bounding box.

    Returns list of dicts: {name, x, y, cx, cy} in EMU (slide-absolute).
    Strategy: walk slide.part.rels for the diagramDrawing part (Microsoft cache that
    contains the rendered SmartArt with explicit per-shape geometry), iterate its <dsp:sp>
    elements, follow each shape's <a:blip r:embed=...> to the SVG image, parse the SVG id
    to get the icon name, and compute slide-absolute coords by adding the frame offset.
    """
    NS_DSP = 'http://schemas.microsoft.com/office/drawing/2008/diagram'
    DRAWING_RELTYPE = 'http://schemas.microsoft.com/office/2007/relationships/diagramDrawing'
    try:
        # Find the diagramDrawing part referenced by this slide (PowerPoint caches it
        # under the slide's rels, not the diagramData's rels)
        drawing_part = None
        for rel in slide_part.rels.values():
            if rel.reltype == DRAWING_RELTYPE:
                drawing_part = rel.target_part
                break
        if drawing_part is None:
            return []
        # Build rId -> svg-id lookup for this drawing part
        rid_to_svg_name = {}
        for rel in drawing_part.rels.values():
            if not rel.reltype.endswith('/image'):
                continue
            tp = rel.target_part
            if not str(getattr(tp, 'partname', '')).lower().endswith('.svg'):
                continue
            try:
                text = tp.blob.decode('utf-8', errors='ignore')
            except Exception:
                continue
            m = re.search(r'<svg[^>]*\bid="([^"]+)"', text)
            if m:
                rid_to_svg_name[rel.rId] = _camel_to_words(m.group(1))
        if not rid_to_svg_name:
            return []
        from lxml import etree
        root = etree.fromstring(drawing_part.blob)
        out = []
        f_left, f_top = frame.left, frame.top
        for sp in root.iter(f'{{{NS_DSP}}}sp'):
            # Prefer <asvg:svgBlip r:embed=rId> (points at the SVG which carries the icon name);
            # fall back to <a:blip> (points at the PNG, which has no name) for shapes without SVG.
            embed_rid = None
            for el in sp.iter():
                ltag = el.tag.split('}')[-1] if '}' in el.tag else el.tag
                if ltag == 'svgBlip':
                    for k, v in el.attrib.items():
                        if k.endswith('}embed') or k == 'embed':
                            embed_rid = v
                            break
                    if embed_rid:
                        break
            if not embed_rid:
                for blip in sp.iter(f'{NS_A}blip'):
                    for k, v in blip.attrib.items():
                        if k.endswith('}embed') or k == 'embed':
                            embed_rid = v
                            break
                    if embed_rid:
                        break
            if not embed_rid or embed_rid not in rid_to_svg_name:
                continue
            # Find the xfrm geometry
            off_x = off_y = ext_cx = ext_cy = None
            for xfrm in sp.iter():
                tag = xfrm.tag.split('}')[-1]
                if tag != 'xfrm':
                    continue
                for child in xfrm:
                    ctag = child.tag.split('}')[-1]
                    if ctag == 'off':
                        off_x = int(child.get('x', 0))
                        off_y = int(child.get('y', 0))
                    elif ctag == 'ext':
                        ext_cx = int(child.get('cx', 0))
                        ext_cy = int(child.get('cy', 0))
                if off_x is not None and ext_cx is not None:
                    break
            if off_x is None or ext_cx is None:
                continue
            out.append({
                'name': rid_to_svg_name[embed_rid],
                'x': f_left + off_x,
                'y': f_top + off_y,
                'cx': ext_cx,
                'cy': ext_cy,
            })
        return out
    except Exception:
        return []


def remove_previous_smartart_icon_caption_shapes(slide):
    spTree = slide.shapes._spTree
    to_remove = []
    for sh in slide.shapes:
        try:
            name = sh.name or ''
        except Exception:
            continue
        if name.startswith(SMARTART_ICON_SHAPE_NAME_PREFIX):
            to_remove.append(sh._element)
    for el in to_remove:
        spTree.remove(el)
    return len(to_remove)


def extract_smartart_text(slide_part, frame):
    """Resolve the SmartArt's diagramData relationship and return list of visible texts.

    Returns [] if the diagram data part cannot be resolved or contains no text.
    """
    # Find the rId of the diagramData relationship referenced by this graphicFrame
    rid = None
    for child in frame._element.iter():
        # dgm:relIds element has r:dm attribute pointing at diagramData rId
        if child.tag.endswith('}relIds'):
            for k, v in child.attrib.items():
                if k.endswith('}dm') or k == 'dm':
                    rid = v
                    break
            if rid:
                break
    if not rid:
        return []
    try:
        rel = slide_part.rels[rid]
        if DIAGRAM_DATA_URI not in rel.reltype:
            return []
        data_part = rel.target_part
        from lxml import etree
        root = etree.fromstring(data_part.blob)
        texts = []
        for t in root.iter(f'{NS_A}t'):
            if t.text and t.text.strip():
                texts.append(t.text.strip())
        # Dedupe consecutive duplicates while preserving order
        out = []
        for t in texts:
            if not out or out[-1] != t:
                out.append(t)
        return out
    except Exception:
        return []


def generate_smartart_caption(texts, icons=None):
    """Build a deterministic accessibility caption from SmartArt content + icon names."""
    icons = icons or []
    if not texts and not icons:
        return None
    parts = []
    if icons:
        parts.append(f"Diagram ({', '.join(icons)} icons):")
    else:
        parts.append("Diagram:")
    if texts:
        parts.append("; ".join(texts) + ".")
    cap = " ".join(parts)
    if len(cap) > 180:
        cap = cap[:177].rstrip() + '...'
    return cap


def remove_previous_smartart_caption_shapes(slide):
    """Remove SmartArt caption text boxes previously added by /captioner."""
    spTree = slide.shapes._spTree
    to_remove = []
    for sh in slide.shapes:
        try:
            name = sh.name or ''
        except Exception:
            continue
        if name.startswith(SMARTART_SHAPE_NAME_PREFIX):
            to_remove.append(sh._element)
    for el in to_remove:
        spTree.remove(el)
    return len(to_remove)


def iter_pictures_recursive(shapes, depth=0):
    """Yield (shape, depth) for every Picture; recurse into GROUP shapes."""
    for sh in shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield sh, depth
            elif sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    yield from iter_pictures_recursive(sh.shapes, depth + 1)
                except NotImplementedError:
                    continue
        except Exception:
            continue


def remove_previous_caption_shapes(slide, target_pic_hash=None,
                                   oracle=None, force_overwrite_edits=False):
    """Remove caption text boxes previously added by /captioner.

    If target_pic_hash given, only remove captions matching that picture's hash
    (encoded in the shape name suffix). Else remove all captioner captions on slide.

    Edit-aware (v0.2.2): if `oracle` (hash -> original caption text from the prior
    audit CSV) is given and a caption's current text DIFFERS from the original, the
    caption is treated as an instructor hand-edit and is PRESERVED (not removed),
    unless force_overwrite_edits=True. Returns (n_removed, n_preserved_edits).
    """
    spTree = slide.shapes._spTree
    to_remove = []
    preserved = 0
    for sh in slide.shapes:
        name = ''
        try:
            name = sh.name or ''
        except Exception:
            continue
        if not name.startswith(CAPTION_SHAPE_NAME_PREFIX):
            continue
        if target_pic_hash is not None and not name.endswith(target_pic_hash):
            continue
        if oracle and not force_overwrite_edits:
            h = name[len(CAPTION_SHAPE_NAME_PREFIX):]
            orig = oracle.get(h)
            try:
                cur = sh.text_frame.text if getattr(sh, 'has_text_frame', False) else ''
            except Exception:
                cur = ''
            if orig is not None and cur.strip() and cur.strip() != orig.strip():
                preserved += 1   # instructor edit — leave it in place
                continue
        to_remove.append(sh._element)
    for el in to_remove:
        spTree.remove(el)
    return len(to_remove), preserved


def load_caption_oracle(audit_dir, deck):
    """hash -> original caption text, from the prior <deck>_audit.csv (if any).
    Used by edit-aware --update-existing to detect instructor hand-edits."""
    oracle = {}
    p = os.path.join(audit_dir, f"{deck}_audit.csv")
    if not os.path.exists(p):
        return oracle
    try:
        with open(p, encoding='utf-8', newline='') as f:
            for r in csv.DictReader(f):
                h, cap = r.get('image_hash', ''), r.get('caption', '')
                if h and cap and cap.strip().lower() not in ('[decorative]', 'decorative'):
                    oracle.setdefault(h, cap)
    except Exception:
        pass
    return oracle


def apply_to_deck(deck_info, captions, captioned_dir, audit_dir, style, opts):
    deck = deck_info['deck']
    src = deck_info['deck_path']

    # Re-run protection: warn if source already looks like a captioner output
    src_base = os.path.basename(src)
    if src_base.endswith('_captioned.pptx'):
        print(f"  WARN: input {src_base} already has _captioned suffix; output will be {src_base[:-5]}_captioned.pptx")

    dst_name = f"{deck}_captioned.pptx" if not deck.endswith('_captioned') else f"{deck}.pptx"
    dst = os.path.join(captioned_dir, dst_name)

    # Edit-aware oracle: prior audit CSV (hash -> original caption). Loaded
    # BEFORE we overwrite the audit, so --update-existing can spot hand-edits.
    caption_oracle = (load_caption_oracle(audit_dir, deck)
                      if opts['update_existing'] else {})

    if opts['dry_run']:
        # Don't copy; we'll only inspect
        prs = Presentation(src)
    else:
        # Back up an existing output before overwriting it (re-run safety).
        if opts['update_existing'] and os.path.exists(dst):
            try:
                shutil.copy(dst, dst + '.bak.pptx')
            except Exception:
                pass
        shutil.copy(src, dst)
        prs = Presentation(dst)

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    audit_rows = []
    n_caption_shapes_removed = 0
    n_preserved_edits = 0
    spell_rows = []
    spell_seen = set()
    qc_rows = []
    qc_seen = set()
    sp_engine, sp_wl = opts.get('_sp'), opts.get('_wl', set())

    # v0.2.5: repeated-background detection. An image whose identical hash
    # appears as a structural picture on >= bg_repeat_threshold distinct slides
    # is template/background chrome (e.g. a slide-wide brick-wall texture). Per-
    # image vision sometimes mislabels these as informative; captioning the same
    # background on every slide is clutter, not accessibility. For such hashes we
    # override a NON-decorative caption to a decorative skip. Hashes the vision
    # pass already marked [decorative] are left untouched (it got those right).
    _bg_thr = opts.get('bg_repeat_threshold', 4)
    _hash_slides = {}
    for _p in deck_info.get('pictures', []):
        _hash_slides.setdefault(_p.get('image_hash'), set()).add(_p.get('slide'))
    bg_hashes = {h for h, sl in _hash_slides.items()
                 if h and _bg_thr and len(sl) >= _bg_thr}

    for s_idx, slide in enumerate(prs.slides, 1):
        # v0.2.1: per-slide placement geometry (shared by SmartArt-icon and
        # main caption paths). Must be defined BEFORE the SmartArt loop so
        # icon captions inherit the same footer clearance + obstacle list.
        footer_limit = _g_slide_footer_top(slide, slide_h)
        effective_footer_limit = footer_limit - FOOTER_CLEARANCE_EMU
        # v0.2.3: EVERY non-empty text frame is a hard obstacle (title, body,
        # AND plain text boxes / auto-shapes — the question/answer/label text the
        # v0.2.2 model was blind to). A caption must never cover any of these.
        text_rects = _g_text_rects(slide)
        # v0.2.2: collect every picture's full (L,T,R,B) rect so a caption
        # for picture-N never lands inside picture-M on the same slide. The
        # caption's own picture is filtered out at use-time so it can still be
        # placed adjacent to its own picture. (Upgraded from y-band to 2D so a
        # caption beside a picture isn't falsely blocked.)
        all_pic_rects: list[tuple[int, int, int, int, str]] = []  # (L,T,R,B,pic_id)
        for _pic_scan in iter_slide_pics(slide):
            _o_left = _pic_scan.get('off_x')
            _o_top = _pic_scan.get('off_y')
            _o_w   = _pic_scan.get('ext_cx')
            _o_h   = _pic_scan.get('ext_cy')
            if None in (_o_left, _o_top, _o_w, _o_h):
                continue
            # v0.2.5: a repeated-background image (see bg_hashes) is NOT a
            # placement obstacle. A full-bleed background texture covers the whole
            # slide, so leaving it in the obstacle set makes every candidate look
            # ">50% inside a picture" and forces flagged-no-slot for the REAL
            # photos on that slide. We don't caption the background, and a caption
            # card sitting over it is fine — so exclude it from obstacles. Real
            # pictures, titles, body text and footers still block captions.
            _rid = _pic_scan.get('rid')
            if _rid is not None and bg_hashes:
                try:
                    if hashlib.sha256(resolve_blob(slide, _rid)).hexdigest()[:12] in bg_hashes:
                        continue
                except Exception:
                    pass
            all_pic_rects.append((int(_o_left), int(_o_top),
                                  int(_o_left) + int(_o_w), int(_o_top) + int(_o_h),
                                  _pic_scan.get('pic_id', '')))
        # Per-slide accumulator for Fix-E (caption-caption overlap avoidance).
        # Stores full (L,T,R,B) rects so clear_all_obstacles_2d reads them
        # correctly (v0.2.1 stored (l,t,w,h), which the (L,T,R,B) obstacle test
        # silently misread as a degenerate band — disabling cap-cap avoidance).
        # MUST reset here, at the top of each slide loop — not before, not inside pic loop.
        placed_caps: list[tuple[int, int, int, int]] = []
        if opts['spellcheck'] or opts['dateqc']:
            for body_text in iter_slide_body_text(slide.shapes):
                if opts['spellcheck']:
                    spell_scan(body_text, 'slide-text', s_idx, sp_engine, sp_wl,
                               spell_seen, spell_rows)
                if opts['dateqc']:
                    qc_scan(body_text, 'slide-text', s_idx, qc_seen, qc_rows)
        # Idempotency: if --update-existing AND not dry-run, strip any prior captioner shapes
        # on this slide before adding new ones.
        if opts['update_existing'] and not opts['dry_run']:
            _rm, _pres = remove_previous_caption_shapes(
                slide, oracle=caption_oracle,
                force_overwrite_edits=opts.get('force_overwrite_edits', False))
            n_caption_shapes_removed += _rm
            n_preserved_edits += _pres
            n_caption_shapes_removed += remove_previous_smartart_caption_shapes(slide)
            n_caption_shapes_removed += remove_previous_smartart_icon_caption_shapes(slide)

        # SmartArt frames: caption deterministically from extracted text content
        if opts['caption_smartart']:
            for sa_idx, frame in enumerate(iter_smartart_frames(slide)):
                placements = extract_smartart_icon_placements(slide.part, frame)
                # Per-icon caption boxes: short label directly under each icon.
                # Text-only SmartArts (no icons) get NO caption — the visible text is
                # already accessible via the slide's own text layer.
                if placements and not opts['dry_run']:
                    icon_h = 200000  # ~0.22"
                    icon_gap = 20000  # tiny gap below icon
                    for ic_idx, p in enumerate(placements):
                        # v0.2.2: apply Fix-B (footer clearance) + Fix-D (widening)
                        # to the SmartArt-icon caption path so it inherits the same
                        # placement-quality guarantees as the main caption path.
                        # Fix-D FIRST: width determines wrap, which determines the
                        # box's RENDERED height (auto_size grows it). Compute width,
                        # then estimate height, so the footer math below uses the
                        # height the box will actually occupy — not the nominal
                        # 200k (a long label wrapping to 3 lines near the footer
                        # otherwise grows down past it: a near-footer icon-caption case).
                        name_text = p['name'] or ''
                        # v0.2.5: icon-proportional sizing. A small icon gets a
                        # smaller caption so the card never dwarfs it. Font scales
                        # 6–8pt with icon width; the width floor, per-char estimate,
                        # per-line height, AND required_caption_width all track the
                        # chosen font so the text still fits and the verify overflow
                        # gate (which reads the real font size) agrees.
                        if p['cx'] < 400000:        # < ~0.44" (tiny icon-strip clip art)
                            ic_font_pt = 5
                        elif p['cx'] < 650000:      # < ~0.71"
                            ic_font_pt = 6
                        elif p['cx'] < 900000:      # < ~0.98"
                            ic_font_pt = 7
                        else:
                            ic_font_pt = 8
                        _emu_per_char = max(1, int(EMU_PER_CHAR_DEFAULT * ic_font_pt / 8))
                        # Small floor only — the box should hug the icon/text, not a
                        # fixed minimum. required_caption_width (longest word) below
                        # widens it whenever the text actually needs more, so a short
                        # label like "new" gets a tight box instead of 0.33".
                        _w_floor = int(150000 * ic_font_pt / 8)
                        ic_width = p['cx'] if p['cx'] >= _w_floor else _w_floor
                        chars_per_line = max(6, ic_width // _emu_per_char)
                        if name_text and len(name_text) > 2 * chars_per_line:
                            # 2-line target. Floor scales with the icon (the full
                            # MIN_CAPTION_WIDTH_EMU floor is for body captions and
                            # would give a tiny icon a huge card).
                            needed = max(_w_floor,
                                         (len(name_text) // 2 + 1) * _emu_per_char)
                            ic_width = min(slide_w - p['x'], needed)
                        # v0.2.4: a single WORD cannot wrap, so the box must be at
                        # least as wide as the longest word (the slide-8 "workflow"
                        # / "electrician" overflow). Widen + RE-CENTER on the icon.
                        _need_w = required_caption_width(name_text, is_icon=True, font_pt=ic_font_pt)
                        if _need_w > ic_width:
                            ic_width = min(slide_w, _need_w)
                        # v0.2.5: cap the box so a tiny icon never gets an oversized
                        # card — at most ~3× the icon width (min ~0.5"), but never
                        # narrower than the longest word (auto_size grows height to
                        # wrap a multi-word label instead of ballooning the width).
                        # ~3x icon width keeps multi-word labels short enough that
                        # they don't wrap tall and collide on dense, near-footer icon
                        # strips (the dense near-footer icon strip). The 5pt font
                        # above is what makes the card smaller; over-narrowing width
                        # just trades width for height and reintroduces overlaps.
                        _w_cap = max(_need_w, min(ic_width, max(p['cx'] * 3, 460000)))
                        ic_width = _w_cap
                        _icon_cx = p['x'] + p['cx'] // 2
                        ic_left = max(0, min(_icon_cx - ic_width // 2, slide_w - ic_width))
                        # Estimate rendered height from wrap (≈160k EMU/line at 8pt,
                        # scaled by font + box margins); never less than a scaled
                        # nominal floor.
                        eff_cpl = max(6, ic_width // _emu_per_char)
                        est_lines = max(1, (len(name_text) + eff_cpl - 1) // eff_cpl)
                        # Tight height: hug the wrapped text (~1.25x line spacing per
                        # line + small top/bottom margins) so the card never extends
                        # below the text. auto_size (SHAPE_TO_FIT_TEXT) also re-fits
                        # in PowerPoint; this keeps the STORED height correct too.
                        _line_emu = int(ic_font_pt * 1.25 * 12700)
                        ic_h = est_lines * _line_emu + 12000  # +2x 6000 EMU margins
                        # Try below the icon first; if that would intrude the footer
                        # clearance band, try above; clamp on-slide as last resort.
                        below_ic_top = p['y'] + p['cy'] + icon_gap
                        above_ic_top = max(0, p['y'] - icon_gap - ic_h)
                        if below_ic_top + ic_h <= effective_footer_limit:
                            ic_top = below_ic_top
                        elif above_ic_top + ic_h <= p['y']:  # fits above icon
                            ic_top = above_ic_top
                        else:
                            ic_top = max(0, effective_footer_limit - ic_h)
                        # v0.2.2: guarantee footer safety regardless of branch.
                        # An icon that itself sits in/below the footer band (a
                        # SmartArt overflowing into the footer zone — the a footer-zone
                        # icon case) would otherwise push its 'above'
                        # caption into the band. Clamp the bottom to the limit.
                        ic_top = min(ic_top, max(0, effective_footer_limit - ic_h))
                        # v0.2.2 Fix-E for SmartArt icons: nudge the icon caption
                        # DOWN past any already-placed caption it would overlap
                        # (the v0.2.1 path placed every icon caption blindly,
                        # producing stacks of 6+ overlapping boxes on icon-heavy
                        # SmartArts). Cap the nudges so we never push past footer.
                        ic_rect = (ic_left, ic_top, ic_left + ic_width, ic_top + ic_h)
                        _tries = 0
                        while (_tries < 8
                               and not clear_all_obstacles_2d(ic_rect, placed_caps, vfrac=0.02)
                               and (ic_top + ic_h + icon_gap + ic_h)
                                   <= effective_footer_limit):
                            ic_top = ic_top + ic_h + icon_gap
                            ic_rect = (ic_left, ic_top, ic_left + ic_width, ic_top + ic_h)
                            _tries += 1
                        # If it STILL overlaps another caption and can't move clear
                        # (e.g. a 6-icon strip whose last icon sits in the footer
                        # zone, so its caption is clamped up into its neighbour —
                        # the dense near-footer icon strip), SKIP rather than ship
                        # an overlap. Parity with the main-caption no-slot policy.
                        if not clear_all_obstacles_2d(ic_rect, placed_caps, vfrac=0.02):
                            audit_rows.append({
                                'slide': s_idx, 'pic_id': '',
                                'image_hash': f'smartart_{sa_idx}_icon_{ic_idx}',
                                'caption': p['name'], 'char_len': len(p['name']),
                                'in_group_depth': 0, 'action': 'flagged-no-slot',
                            })
                            continue
                        # Track this icon caption as an obstacle (full L,T,R,B
                        # rect) for subsequent icon + main-caption placements.
                        placed_caps.append((ic_left, ic_top, ic_left + ic_width, ic_top + ic_h))
                        itb = slide.shapes.add_textbox(ic_left, ic_top, ic_width, ic_h)
                        try:
                            itb._element.nvSpPr.cNvPr.set('name', f"{SMARTART_ICON_SHAPE_NAME_PREFIX}{sa_idx}_{ic_idx}")
                        except Exception:
                            pass
                        if style['bg_color']:
                            itb.fill.solid()
                            itb.fill.fore_color.rgb = RGBColor.from_string(style['bg_color'])
                        if style['border_color']:
                            itb.line.color.rgb = RGBColor.from_string(style['border_color'])
                            itb.line.width = Emu(6350)
                        itf = itb.text_frame
                        itf.word_wrap = True
                        # Auto-grow vertical to fit wrapped text (icons can be narrow)
                        itf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                        itf.margin_left = Emu(20000); itf.margin_right = Emu(20000)
                        itf.margin_top = Emu(6000); itf.margin_bottom = Emu(6000)
                        ip = itf.paragraphs[0]
                        ip.alignment = PP_ALIGN.CENTER
                        irun = ip.add_run()
                        irun.text = p['name']
                        irun.font.size = Pt(ic_font_pt)
                        irun.font.italic = style['italic']
                        irun.font.name = style['font_name']
                        irun.font.color.rgb = RGBColor.from_string(style['font_color'])
                        audit_rows.append({
                            'slide': s_idx, 'pic_id': '', 'image_hash': f'smartart_{sa_idx}_icon_{ic_idx}',
                            'caption': p['name'], 'char_len': len(p['name']), 'in_group_depth': 0,
                            'action': 'added-smartart-icon',
                        })
                # Note: SmartArts with NO icons (text-only diagrams) get NO caption.
                # The visible text already lives in the slide's accessible text layer.
                if not placements:
                    audit_rows.append({
                        'slide': s_idx, 'pic_id': '', 'image_hash': f'smartart_{sa_idx}',
                        'caption': '', 'char_len': 0, 'in_group_depth': 0,
                        'action': 'skipped-smartart-text-only',
                    })

        # (per-slide geometry init moved above the SmartArt loop)
        for pic in iter_slide_pics(slide):
            depth = pic['depth']
            pic_id = pic['pic_id']
            if pic['rid'] is None:
                # Linked/external pic (r:link, no r:embed) or malformed blip:
                # cannot hash bytes -> skip + structured audit row (parity
                # with extract_images.py, which also skips + logs these).
                audit_rows.append({
                    'slide': s_idx, 'pic_id': pic_id, 'image_hash': '',
                    'caption': '', 'char_len': 0, 'in_group_depth': depth,
                    'action': 'skipped-linked-no-embed',
                })
                continue
            try:
                blob = resolve_blob(slide, pic['rid'])
                ext = guess_ext(slide.part.related_part(pic['rid']))
            except Exception:
                audit_rows.append({
                    'slide': s_idx, 'pic_id': pic_id, 'image_hash': '',
                    'caption': '', 'char_len': 0, 'in_group_depth': depth,
                    'action': 'skipped-image-extract-failed',
                })
                continue
            h = hashlib.sha256(blob).hexdigest()[:12]
            key = f"{deck}/{h}.{ext}"
            caption = captions.get(key)

            if caption is None:
                audit_rows.append({
                    'slide': s_idx, 'pic_id': pic_id, 'image_hash': h,
                    'caption': '', 'char_len': 0, 'in_group_depth': depth,
                    'action': 'skipped-no-caption',
                })
                continue
            if caption.strip().lower() in ('[decorative]', 'decorative', ''):
                audit_rows.append({
                    'slide': s_idx, 'pic_id': pic_id, 'image_hash': h,
                    'caption': '[decorative]', 'char_len': 0, 'in_group_depth': depth,
                    'action': 'skipped-decorative',
                })
                continue
            if h in bg_hashes:
                # Repeated template/background image (e.g. a slide-wide texture
                # on every slide). Vision gave it a caption, but captioning the
                # same background N times is clutter — skip as decorative.
                audit_rows.append({
                    'slide': s_idx, 'pic_id': pic_id, 'image_hash': h,
                    'caption': '[repeated-background]', 'char_len': 0,
                    'in_group_depth': depth,
                    'action': 'skipped-decorative-background',
                })
                continue

            if opts['spellcheck']:
                spell_scan(caption, 'caption', s_idx, sp_engine, sp_wl,
                           spell_seen, spell_rows)
            if opts['dateqc']:
                qc_scan(caption, 'caption', s_idx, qc_seen, qc_rows)

            # Compute placement. Raw <a:off>/<a:ext> EMU — proven byte-identical
            # to the old python-pptx pic.left/.top/.width/.height for every deck
            # pic in this corpus (incl. group-nested), so caption cards do not
            # move. Defensive 0 default if a pic somehow lacks an xfrm (the
            # off-slide fallback-placement logic below then handles it).
            p_left = pic['off_x']
            p_top = pic['off_y']
            p_width = pic['ext_cx']
            p_height = pic['ext_cy']
            # Placeholder-hosted picture: geometry inherited from the layout.
            # Resolve it via the python-pptx placeholder, else the caption
            # lands off-page at (0, 50000) (the 2026-05-18 slide-3 bug).
            if (None in (p_left, p_top, p_width, p_height)
                    and pic.get('is_placeholder') and pic.get('ph_idx') is not None):
                g = resolve_ph_geometry(slide, pic['ph_idx'])
                if g is not None:
                    p_left, p_top, p_width, p_height = g
            # Last-resort: center a default box on the slide rather than the
            # old degenerate top-left stub.
            if None in (p_left, p_top, p_width, p_height):
                p_width = p_width or min(4000000, slide_w)
                p_height = p_height or 0
                p_left = p_left if p_left is not None else (slide_w - p_width) // 2
                p_top = p_top if p_top is not None else int(slide_h * 0.55)

            gap = opts['gap_emu']
            # --- Width (Fix-D): default = picture width, widen if the caption
            # wouldn't fit in ≤2 lines at that width. Never truncate. ---
            c_left = max(0, min(p_left, slide_w - 500000))
            base_width = min(p_width, slide_w - c_left)
            if base_width < 500000:
                base_width = min(500000, slide_w)
            est_chars_per_line_at_base = max(10, base_width // EMU_PER_CHAR_DEFAULT)
            if len(caption) > 2 * est_chars_per_line_at_base:
                needed = max(MIN_CAPTION_WIDTH_EMU,
                             (len(caption) // 2 + 1) * EMU_PER_CHAR_DEFAULT)
                c_width = min(slide_w - c_left, needed)
            else:
                c_width = base_width
            # v0.2.4: a single WORD cannot wrap — the box must be at least as wide
            # as the longest word or it overflows horizontally (the slide-8 defect).
            _need_w = required_caption_width(caption, is_icon=False)
            if _need_w > c_width:
                c_width = min(slide_w, _need_w)
            if c_left + c_width > slide_w:
                c_left = max(0, slide_w - c_width)

            # --- Rendered height estimate (v0.2.3). The caption box auto-sizes,
            # so a 2-line caption occupies MORE than the nominal height. Use the
            # estimate for every placement/overlap decision so a grown box never
            # spills onto a neighbor (the slide-8 "label falls onto next" defect). ---
            c_height = estimate_caption_height(caption, c_width,
                                               nominal_emu=opts['height_emu'])

            # --- Obstacles (v0.2.3). HARD obstacles a caption must NEVER cover:
            # every text frame on the slide (title, body, AND plain text boxes /
            # auto-shapes — the v0.2.2 model only saw placeholders) plus every
            # OTHER picture. The caption's own picture is NOT a hard obstacle: a
            # clean slot sits just outside it, and its bottom strip is an allowed
            # last-resort band (covering a little image beats covering text). ---
            # Exclude ONLY the caption's own picture from the obstacle set —
            # matched by GEOMETRY, not pic_id. iter_slide_pics assigns the same
            # pic_id (often 0) to multiple pictures on a slide, so a pic_id match
            # wrongly drops a DIFFERENT picture from the obstacles and lets the
            # caption land on it (two stacked pictures sharing a pic_id). Geometry is
            # unique per visible picture; if two pictures truly coincide, excluding
            # both is harmless (they occupy the same space).
            own_pic = (p_left, p_top, p_left + p_width, p_top + p_height)
            other_pics = [(_l, _t, _r, _b) for _l, _t, _r, _b, _pid in all_pic_rects
                          if (_l, _t, _r, _b) != own_pic]
            # Two tolerances: a caption must BARELY touch text, let alone cover it
            # (TEXT_VFRAC, near-zero) — but a small overlap with a picture's edge is
            # visually fine (PIC_VFRAC). Clean candidates avoid even the own picture;
            # the band (with_own_pic=False) deliberately sits in its own picture.
            TEXT_VFRAC = 0.03
            PIC_VFRAC = 0.15
            CAP_VFRAC = 0.02

            def _clears2(left, top, width, height, with_own_pic=True):
                rect = (left, top, left + width, top + height)
                pics = other_pics + ([own_pic] if with_own_pic else [])
                return (clear_all_obstacles_2d(rect, text_rects, vfrac=TEXT_VFRAC)
                        and clear_all_obstacles_2d(rect, pics, vfrac=PIC_VFRAC)
                        and clear_all_obstacles_2d(rect, placed_caps, vfrac=CAP_VFRAC))

            def _h_for(w):
                return estimate_caption_height(caption, w, nominal_emu=opts['height_emu'])

            # --- Clean external candidates: below -> above -> RIGHT-of-pic ->
            # LEFT-of-pic. v0.2.4 adds the side candidates so a structural picture
            # (numbered chevron / icon) gets a caption BESIDE it instead of a band
            # burying its content. Each is (action, left, top, width). ---
            below_top = p_top + p_height + gap
            above_top = p_top - gap - c_height
            _vc = p_top + max(0, (p_height - c_height) // 2)   # vert-center for side boxes
            pick = None  # (action, left, top, width, height)
            # 1. Below the picture (external, cleanest — no image overlap at all).
            if (below_top + c_height <= effective_footer_limit
                    and _clears2(c_left, below_top, c_width, _h_for(c_width))):
                pick = ('added', c_left, below_top, c_width, _h_for(c_width))
            # 2. Small bottom-of-picture band — PREFERRED over above/beside so the
            #    caption stays WITH its own photo (at its bottom) instead of drifting
            #    up into the title/body text region (a text-crowded placement). The band is
            #    small (BAND_FONT_PT). Skipped for a small/thin/structural picture
            #    where a band would bury its content (it then falls to above/beside).
            if pick is None:
                band_w = min(c_width, p_width) if (p_width and p_width >= 500000) else c_width
                _band_cw = max(1, int(BAND_FONT_PT * 0.5 * 12700))      # 0.5em/char
                _band_cpl = max(6, band_w // _band_cw)
                _band_lines = max(1, (len(caption) + _band_cpl - 1) // _band_cpl)
                band_h = int(_band_lines * BAND_FONT_PT * 1.3 * 12700) + 24000
                band_left = max(0, min(p_left, slide_w - band_w))
                band_top = max(0, min(p_top + p_height - band_h,
                                      effective_footer_limit - band_h))
                band_rect = (band_left, band_top, band_left + band_w, band_top + band_h)
                if (band_left + band_w <= slide_w
                        and not band_covers_structural_picture(band_rect, own_pic)
                        and _clears2(band_left, band_top, band_w, band_h, with_own_pic=False)):
                    pick = ('inside-bottom', band_left, band_top, band_w, band_h)
            # 3. Above / beside the picture — last resort before skipping.
            if pick is None:
                cands = []
                if above_top >= 0:
                    cands.append(('fallback-above', c_left, above_top, c_width))
                _r_left = p_left + p_width + gap
                _r_w = min(c_width, slide_w - _r_left)
                if _r_w >= 500000:
                    cands.append(('side-right', _r_left, _vc, _r_w))
                _l_w = min(c_width, p_left - gap)
                if _l_w >= 500000:
                    cands.append(('side-left', max(0, p_left - gap - _l_w), _vc, _l_w))
                for action_name, left, top, w in cands:
                    if _clears2(left, top, w, _h_for(w)):
                        pick = (action_name, left, top, w, _h_for(w)); break
                if pick is None:
                    # Horizontal nudge of the above candidate.
                    for action_name, left, top, w in cands:
                        if action_name != 'fallback-above':
                            continue
                        for step in (w + gap, (w + gap) // 2 + gap):
                            alt = left + step
                            if alt + w <= slide_w and _clears2(alt, top, w, _h_for(w)):
                                pick = (action_name, alt, top, w, _h_for(w)); break
                        if pick is not None:
                            break
            if pick is None:
                # Truly no slot that avoids text — SKIP + flag (never cover text).
                audit_rows.append({
                    'slide': s_idx, 'pic_id': pic_id, 'image_hash': h,
                    'caption': caption, 'char_len': len(caption),
                    'in_group_depth': depth,
                    'action': f'{"dry-run-would-" if opts["dry_run"] else ""}'
                              f'flagged-no-slot',
                })
                continue
            action, c_left, c_top, c_width, c_height = pick
            is_band = (action == 'inside-bottom')

            # Final hard clamp — fully on-slide, footer-safe (band keeps its top
            # so it stays bottom-aligned to the picture).
            c_top = max(0, min(c_top, slide_h - c_height))
            if c_top + c_height > slide_h:
                c_height = max(MIN_CAPTION_HEIGHT, slide_h - c_top)
            if not is_band and c_top + c_height > effective_footer_limit:
                c_height = max(MIN_CAPTION_HEIGHT, effective_footer_limit - c_top)

            # === v0.2.4 PRE-WRITE INVARIANT — the standing guarantee ===
            # Re-verify the FINAL clamped caption rect against EVERY defect class
            # right before writing it. If any upstream step (the clamp, the height
            # estimate, a new/unseen layout, group-nested coordinate drift) produced
            # a defect, REFUSE to place the caption — skip + flag — rather than ship
            # it. A placement defect is therefore impossible to WRITE by construction;
            # verify.py / the auditor are defense-in-depth, not the only line.
            _final = (c_left, c_top, c_left + c_width, c_top + c_height)
            _viol = None
            if not (c_left >= 0 and c_top >= 0
                    and c_left + c_width <= slide_w and c_top + c_height <= slide_h):
                _viol = 'off-slide'
            elif c_top + c_height > effective_footer_limit + 1:
                _viol = 'footer'
            elif not clear_all_obstacles_2d(_final, text_rects, vfrac=TEXT_VFRAC):
                _viol = 'text-overlap'
            elif not clear_all_obstacles_2d(_final, placed_caps, vfrac=CAP_VFRAC):
                _viol = 'caption-overlap'
            elif caption_overflows(caption, c_width, is_icon=False):
                _viol = 'overflow'
            else:
                _fa = max(1, c_width * c_height)
                _pics_to_check = other_pics if is_band else (other_pics + [own_pic])
                for _pl, _pt, _pr, _pb in _pics_to_check:
                    _ix = max(0, min(c_left + c_width, _pr) - max(c_left, _pl))
                    _iy = max(0, min(c_top + c_height, _pb) - max(c_top, _pt))
                    if (_ix * _iy) / _fa > 0.50:
                        _viol = 'in-picture'
                        break
            if _viol is not None:
                cov = visible_coverage(p_left, p_top, p_width, p_height, slide_w, slide_h)
                audit_rows.append({
                    'slide': s_idx, 'pic_id': pic_id, 'image_hash': h,
                    'caption': caption, 'char_len': len(caption),
                    'in_group_depth': depth,
                    'action': f'{"dry-run-would-" if opts["dry_run"] else ""}'
                              f'flagged-self-check-{_viol}',
                    'visible_coverage': round(cov, 3),
                })
                continue  # PRE-WRITE INVARIANT: never write a defective caption
            # === end invariant ===

            # Register this caption (full L,T,R,B rect, at its ESTIMATED grown
            # size) as an obstacle for subsequent pictures on the same slide.
            placed_caps.append((c_left, c_top, c_left + c_width, c_top + c_height))

            if opts['dry_run']:
                audit_rows.append({
                    'slide': s_idx, 'pic_id': pic_id, 'image_hash': h,
                    'caption': caption, 'char_len': len(caption), 'in_group_depth': depth,
                    'action': f'dry-run-would-{action}',
                })
                continue

            # Add the text box
            tb = slide.shapes.add_textbox(c_left, c_top, c_width, c_height)
            # Idempotency: tag the shape name with our prefix + image hash. Band
            # captions get the band prefix so the audit/verify in-picture checks
            # know their own-picture overlap is intentional (not a defect).
            name_prefix = CAPTION_BAND_NAME_PREFIX if is_band else CAPTION_SHAPE_NAME_PREFIX
            try:
                tb._element.nvSpPr.cNvPr.set('name', f"{name_prefix}{h}")
            except Exception:
                pass  # if naming fails, the caption still renders correctly

            # Solid background fill — ensures caption is readable on dark slide backgrounds
            if style['bg_color']:
                tb.fill.solid()
                tb.fill.fore_color.rgb = RGBColor.from_string(style['bg_color'])
            # Thin border so the caption "card" reads cleanly on any background
            if style['border_color']:
                tb.line.color.rgb = RGBColor.from_string(style['border_color'])
                tb.line.width = Emu(6350)  # 0.5pt

            tf = tb.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            _m, _mv = (12000, 6000) if is_band else (50000, 20000)
            tf.margin_left = Emu(_m); tf.margin_right = Emu(_m)
            tf.margin_top = Emu(_mv); tf.margin_bottom = Emu(_mv)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = caption
            run.font.size = Pt(BAND_FONT_PT if is_band else style['font_size'])
            run.font.italic = style['italic']
            run.font.name = style['font_name']
            run.font.color.rgb = RGBColor.from_string(style['font_color'])

            audit_rows.append({
                'slide': s_idx, 'pic_id': pic_id, 'image_hash': h,
                'caption': caption, 'char_len': len(caption), 'in_group_depth': depth,
                'action': action,
            })

    if not opts['dry_run']:
        prs.save(dst)
    if audit_rows:
        # Dry-run writes its OWN file so it can never clobber a prior real
        # apply audit (the 2026-05-18 overwrite hiccup).
        audit_name = f"{deck}_audit_dryrun.csv" if opts['dry_run'] else f"{deck}_audit.csv"
        csv_path = os.path.join(audit_dir, audit_name)
        # v0.2.1: union of keys across all rows — different audit-row types
        # (overlay-fullbleed, flagged-no-slot, etc.) carry additional fields.
        all_keys = []
        seen_keys = set()
        for r in audit_rows:
            for k in r.keys():
                if k not in seen_keys:
                    seen_keys.add(k); all_keys.append(k)
        with open(csv_path, 'w', newline='') as f:
            w = csv.DictWriter(f, fieldnames=all_keys, extrasaction='ignore')
            w.writeheader()
            for r in audit_rows:
                w.writerow(r)

    # QC artifacts go to their OWN directory, never the caption audit dir.
    qc_dir = opts['qc_dir']
    if opts['spellcheck'] and spell_rows:
        sc_path = os.path.join(qc_dir, f"{deck}_spellcheck.csv")
        with open(sc_path, 'w', newline='') as f:
            w = csv.DictWriter(f, fieldnames=['slide', 'source', 'term',
                                              'suggestion', 'known_bad',
                                              'verify_name', 'context'])
            w.writeheader()
            for r in spell_rows:
                w.writerow(r)
    if opts['dateqc'] and qc_rows:
        qp = os.path.join(qc_dir, f"{deck}_qc.csv")
        with open(qp, 'w', newline='') as f:
            w = csv.DictWriter(f, fieldnames=['slide', 'source', 'kind',
                                              'detail', 'context'])
            w.writeheader()
            for r in qc_rows:
                w.writerow(r)

    return {
        'deck': deck, 'dst': dst if not opts['dry_run'] else None,
        'rows': audit_rows, 'n_caption_shapes_removed': n_caption_shapes_removed,
        'n_preserved_edits': n_preserved_edits,
        'spell_rows': spell_rows, 'qc_rows': qc_rows,
    }


def main():
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument('work_dir', help='Directory containing manifest.json and captions.json')
    ap.add_argument('--dry-run', action='store_true', help='Compute caption placements + write audit CSV but do NOT modify any .pptx')
    ap.add_argument('--font-name', default='Calibri', help='Caption font (default: Calibri)')
    ap.add_argument('--font-size', type=float, default=10.0, help='Caption font size in pt (default: 10)')
    ap.add_argument('--font-color', default='333333', help='Caption hex color, no # (default: 333333 dark gray)')
    ap.add_argument('--bg-color', default='FFFFFF', help='Caption text box fill hex, no # (default: FFFFFF white; "" to disable for transparent)')
    ap.add_argument('--border-color', default='CCCCCC', help='Caption text box border hex, no # (default: CCCCCC light gray; "" to disable)')
    ap.add_argument('--italic', type=lambda v: v.lower() in ('1','true','yes','y'), default=True, help='Caption italic (default: true)')
    ap.add_argument('--gap-emu', type=int, default=50000, help='EMU between picture and caption (default: 50000 = ~0.05")')
    ap.add_argument('--height-emu', type=int, default=400000, help='Caption text box height in EMU (default: 400000 = ~0.44")')
    ap.add_argument('--update-existing', action='store_true', help='Detect and remove previously-added captioner shapes before adding new ones (else re-runs DUPLICATE). Edit-aware: a caption whose text no longer matches the prior audit CSV is treated as an instructor hand-edit and is PRESERVED, not stripped.')
    ap.add_argument('--force-overwrite-edits', action='store_true', help='With --update-existing, remove even captions that look instructor-edited (default: preserve them).')
    ap.add_argument('--no-smartart', action='store_true', help='Disable SmartArt captioning (SmartArt captions are auto-generated from diagram text content, on by default)')
    ap.add_argument('--bg-repeat-threshold', type=int, default=4, help='An image whose identical hash appears on >= this many slides is treated as a repeated background/template and skipped (decorative), even if vision captioned it. Default 4; set 0 to disable.')
    ap.add_argument('--quiet', action='store_true', help='Suppress per-deck progress lines')
    # QC is ON BY DEFAULT (spell-check + name-verify tagging + date/template QC).
    ap.add_argument('--quick', action='store_true', help='Captioning ONLY — skip ALL QC (spell-check + date/template scan). Use when you explicitly do not want QC.')
    ap.add_argument('--no-spellcheck', action='store_true', help='Disable just the spell-check pass (date/template QC still runs unless --quick).')
    ap.add_argument('--no-dateqc', action='store_true', help='Disable just the date/doubled-word/leftover-template scan (spell-check still runs unless --quick).')
    ap.add_argument('--spellcheck', action='store_true', help=argparse.SUPPRESS)  # back-compat no-op: spell-check is now default-on
    args = ap.parse_args()
    # Default-on QC with granular + master (--quick) toggles.
    qc_on = not args.quick
    do_spell = qc_on and not args.no_spellcheck
    do_dateqc = qc_on and not args.no_dateqc

    work = os.path.abspath(args.work_dir)
    with open(os.path.join(work, 'manifest.json')) as f:
        manifest = json.load(f)
    captions_path = os.path.join(work, 'captions.json')
    if not os.path.exists(captions_path):
        print(f"ERROR: {captions_path} not found. Write it after reading images via Read tool.")
        sys.exit(2)
    with open(captions_path) as f:
        captions = json.load(f)

    captioned_dir = os.path.join(work, 'captioned_decks')
    audit_dir = os.path.join(work, 'audit')
    qc_dir = os.path.join(work, 'qc')
    os.makedirs(captioned_dir, exist_ok=True)
    os.makedirs(audit_dir, exist_ok=True)
    if do_spell or do_dateqc:
        os.makedirs(qc_dir, exist_ok=True)

    style = {
        'font_name': args.font_name, 'font_size': args.font_size,
        'font_color': args.font_color, 'italic': args.italic,
        'bg_color': args.bg_color, 'border_color': args.border_color,
    }
    sp_engine, sp_wl, sp_status = init_spellcheck(do_spell)
    # FAIL LOUD: spell-check is default-on; a missing optional dep must NOT
    # silently no-op (the 2026-05-18 silent-skip hiccup). Captioning still
    # proceeds, but we banner it, drop a marker, and exit non-zero.
    spellcheck_degraded = do_spell and sp_status == 'unavailable'
    opts = {
        'dry_run': args.dry_run, 'update_existing': args.update_existing,
        'force_overwrite_edits': args.force_overwrite_edits,
        'gap_emu': args.gap_emu, 'height_emu': args.height_emu,
        'caption_smartart': not args.no_smartart,
        'bg_repeat_threshold': args.bg_repeat_threshold,
        'spellcheck': do_spell and sp_status == 'ready',
        'dateqc': do_dateqc, 'qc_dir': qc_dir,
        '_sp': sp_engine, '_wl': sp_wl,
    }

    if args.dry_run:
        print("=== DRY RUN — no .pptx files will be modified ===")
    if args.quick:
        print("=== --quick: captioning ONLY, all QC skipped by request ===")
    if spellcheck_degraded:
        os.makedirs(qc_dir, exist_ok=True)
        with open(os.path.join(qc_dir, 'SPELLCHECK_NOT_RUN.txt'), 'w') as mf:
            mf.write("Spell-check is default-ON but pyspellchecker is NOT importable "
                     "in the Python that ran apply_captions.py.\n"
                     "Captions were applied; SPELL-CHECK DID NOT RUN.\n"
                     "Fix: run via a Python with pyspellchecker installed, or pass "
                     "--no-spellcheck / --quick to acknowledge skipping it.\n")
        print("*" * 72)
        print("*** ERROR: spell-check is ON by default but pyspellchecker is NOT")
        print("*** installed in this Python. Captions WILL be applied, but the")
        print("*** spell-check DID NOT RUN. Marker: qc/SPELLCHECK_NOT_RUN.txt")
        print("*** Re-run with pyspellchecker available, or pass --no-spellcheck")
        print("*** / --quick to explicitly proceed without it. (exit code 3)")
        print("*" * 72)
    elif opts['spellcheck']:
        print(f"=== spell-check ON (default; FLAG-ONLY, whitelist={len(sp_wl)} "
              "terms): captions + slide text scanned; qc/<deck>_spellcheck.csv "
              "emitted; NO .pptx edited, NO auto-correction ===")
    if opts['dateqc']:
        print("=== date/template QC ON (default; FLAG-ONLY): "
              "qc/<deck>_qc.csv emitted ===")

    total_added = total_skipped = total_pics = total_removed = total_errored = 0
    total_spell = total_qc = total_verify = total_preserved = 0
    decks_with_errors = []
    action_counter = Counter()   # categorized coverage across all decks
    for i, (deck_name, deck_info) in enumerate(manifest.items(), 1):
        if 'error' in deck_info:
            total_errored += 1
            decks_with_errors.append(deck_name)
            if not args.quiet:
                print(f"[{i}/{len(manifest)}] {deck_name:<35} SKIPPED (extract error: {deck_info['error'][:60]})")
            continue
        try:
            result = apply_to_deck(deck_info, captions, captioned_dir, audit_dir, style, opts)
        except Exception as e:
            total_errored += 1
            decks_with_errors.append(deck_name)
            if not args.quiet:
                print(f"[{i}/{len(manifest)}] {deck_name:<35} ERROR: {type(e).__name__}: {str(e)[:60]}")
            continue
        rows = result['rows']
        action_counter.update(r['action'] for r in rows)
        added = sum(1 for r in rows if r['action'].startswith('added') or r['action'] in ('fallback-above','fallback-bottom') or r['action'].startswith('dry-run-would'))
        skipped = sum(1 for r in rows if r['action'].startswith('skipped'))
        if not args.quiet:
            verb = 'would-add' if args.dry_run else 'added'
            dst_label = os.path.basename(result['dst']) if result['dst'] else '(dry-run, no output)'
            removed_note = f"  removed-prior:{result['n_caption_shapes_removed']}" if result['n_caption_shapes_removed'] else ''
            print(f"[{i}/{len(manifest)}] {deck_name:<35} {added:>3}/{len(rows)} {verb}  {skipped:>3} skip{removed_note}  -> {dst_label}")
        total_added += added; total_skipped += skipped; total_pics += len(rows)
        total_removed += result['n_caption_shapes_removed']
        total_preserved += result.get('n_preserved_edits', 0)
        sr = result.get('spell_rows', [])
        total_spell += len(sr)
        total_verify += sum(1 for r in sr if r.get('verify_name'))
        total_qc += len(result.get('qc_rows', []))

    print(f"\n{'='*72}")
    print(f"TOTAL: {total_added}/{total_pics} captions {'would be added (dry-run)' if args.dry_run else 'added'}, {total_skipped} skipped, {total_errored} deck-level errors")
    if total_removed:
        print(f"Prior captioner shapes removed (idempotency): {total_removed}")
    if total_preserved:
        print(f"Instructor-edited captions PRESERVED (not overwritten): {total_preserved} "
              "— pass --force-overwrite-edits to replace them instead.")
    if opts['spellcheck']:
        print(f"Spell-check flags (review only, no edits made): {total_spell} "
              f"— see qc/<deck>_spellcheck.csv in {qc_dir}")
        if total_verify:
            print(f"  ↳ {total_verify} are likely PROPER NOUNS (verify_name=True): "
                  "web-verify the canonical spelling BEFORE presenting as a fix "
                  "(see SKILL.md “Name verification”).")
    if opts['dateqc']:
        print(f"Date/template QC flags (review only): {total_qc} "
              f"— see qc/<deck>_qc.csv in {qc_dir}")
    if spellcheck_degraded:
        print("SPELL-CHECK DID NOT RUN — pyspellchecker missing "
              "(qc/SPELLCHECK_NOT_RUN.txt). See banner above.")
    if decks_with_errors:
        print(f"Decks with errors: {decks_with_errors}")
    if not args.dry_run:
        print(f"Captioned decks: {captioned_dir}")
    print(f"Audit CSVs:      {audit_dir}")
    if do_spell or do_dateqc:
        print(f"QC CSVs:         {qc_dir}")

    # ---- Categorized coverage report (v0.2.2) --------------------------------
    # Separate a REAL accessibility gap ("no clean slot" — picture left
    # uncaptioned) from a NON-gap (decorative / text-only SmartArt) so the
    # operator can report coverage honestly instead of a single lumped "skipped".
    def _sum(*prefixes_or_exact):
        n = 0
        for act, cnt in action_counter.items():
            a = act.replace('dry-run-would-', '')
            if a in prefixes_or_exact:
                n += cnt
        return n
    placed = _sum('added', 'fallback-above', 'fallback-bottom', 'added-smartart-icon')
    no_slot = _sum('overlay-fullbleed', 'flagged-no-slot',
                   'flagged-self-check-text-overlap', 'flagged-self-check-footer',
                   'flagged-self-check-caption-overlap', 'flagged-self-check-in-picture',
                   'flagged-self-check-off-slide')   # real WCAG gap (incl. invariant skips)
    decorative = _sum('skipped-decorative')                   # not a gap
    decorative_bg = _sum('skipped-decorative-background')     # repeated template/background — not a gap
    text_only_sa = _sum('skipped-smartart-text-only')         # not a gap (text already accessible)
    errored_imgs = _sum('skipped-no-caption', 'skipped-linked-no-embed',
                        'skipped-image-extract-failed')
    cov_path = os.path.join(work, 'coverage_report.csv')
    with open(cov_path, 'w', newline='') as f:
        w = csv.writer(f)
        w.writerow(['category', 'count', 'is_accessibility_gap', 'note'])
        w.writerow(['captions_placed', placed, 'no', 'visible caption added'])
        w.writerow(['no_clean_slot', no_slot, 'YES',
                    'picture left uncaptioned — no clean placement; needs human review'])
        w.writerow(['decorative', decorative, 'no', 'classified decorative — intentionally uncaptioned'])
        w.writerow(['decorative_background', decorative_bg, 'no',
                    f'repeated background/template image (appears on >= bg-repeat-threshold slides) — intentionally uncaptioned'])
        w.writerow(['smartart_text_only', text_only_sa, 'no',
                    'text-only SmartArt — text already in the accessible layer'])
        w.writerow(['extract_or_link_error', errored_imgs, 'maybe',
                    'image could not be read/hashed — inspect'])
        for act, cnt in sorted(action_counter.items()):
            w.writerow([f'raw:{act}', cnt, '', ''])
    print(f"Coverage report: {cov_path}")
    print(f"  captions placed: {placed} | no-clean-slot (REAL gap, flagged): {no_slot} "
          f"| decorative: {decorative} | repeated-background: {decorative_bg} | text-only SmartArt: {text_only_sa}"
          + (f" | extract/link errors: {errored_imgs}" if errored_imgs else ""))
    if no_slot:
        print(f"  ⚠ Report to stakeholders: {no_slot} pictures have NO visible caption "
              "(no clean placement slot) — see overlay-fullbleed/flagged-no-slot audit rows.")
    print(f"{'='*72}")
    if spellcheck_degraded:
        sys.exit(3)


if __name__ == '__main__':
    main()
